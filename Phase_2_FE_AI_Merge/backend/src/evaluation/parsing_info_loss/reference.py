"""Task-specific reference builders from original files or OmniDocBench GT."""

from __future__ import annotations

import hashlib
import html
import json
import re
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple
from xml.etree import ElementTree as ET

from .schemas import Component, DocumentComponents, SectionNode
from .utils import CAPTION_RE, normalize_text, stable_id

W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def build_office_reference(path: Path, doc_id: str, modality: str) -> DocumentComponents:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return build_docx_reference(path, doc_id)
    if suffix in {".xlsx", ".xlsm"}:
        return build_xlsx_reference(path, doc_id)
    if suffix == ".pptx":
        return build_pptx_reference(path, doc_id)
    raise ValueError(f"Unsupported native Office reference type: {path}")


def build_pdf_weak_reference(path: Path, doc_id: str) -> DocumentComponents:
    """Build a weak PDF reference from extractable page text.

    This is not a substitute for OmniDocBench human GT. It exists for local PDFs
    where the original PDF is the only available reference. The unit is page text:
    enough to measure text coverage and coarse reading order, but not enough to
    make strong claims about table structure or figure-caption association.
    """

    doc = DocumentComponents(
        doc_id=doc_id,
        modality="pdf",
        source_path=str(path),
        reference_type="pdf_text_extraction",
        gt_strength="weak",
    )
    pages = _extract_pdf_pages(path)
    comp_idx = 0
    for page_idx, page_text in enumerate(pages, 1):
        cleaned = normalize_text(page_text)
        if not cleaned:
            continue
        comp_idx += 1
        doc.components.append(
            Component(
                id=stable_id("gt", comp_idx),
                type="text",
                text=cleaned,
                order=comp_idx,
                scope_id="document",
                section_path=[],
                meta={"page": page_idx},
            )
        )
    return doc


def _extract_pdf_pages(path: Path) -> List[str]:
    try:
        import fitz  # PyMuPDF

        with fitz.open(path) as pdf:
            return [page.get_text("text") or "" for page in pdf]
    except Exception:
        pass
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            return [page.extract_text() or "" for page in pdf.pages]
    except Exception:
        pass
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:
        raise RuntimeError(f"Could not extract text from PDF reference: {path}") from exc


def build_docx_reference(path: Path, doc_id: str) -> DocumentComponents:
    doc = DocumentComponents(
        doc_id=doc_id,
        modality="docx",
        source_path=str(path),
        reference_type="native_docx_xml",
        gt_strength="strong",
    )
    with zipfile.ZipFile(path) as zf:
        document_xml = zf.read("word/document.xml")
        rels = _read_rels(zf, "word/_rels/document.xml.rels")
        media_hashes = _media_hashes(zf, "word/media/")
    root = ET.fromstring(document_xml)
    body = root.find(W + "body")
    if body is None:
        return doc

    order = 0
    section_stack: List[SectionNode] = []
    section_counter = 0
    component_counter = 0

    def next_order() -> int:
        nonlocal order
        order += 1
        return order

    def add_component(component_type: str, **kwargs: Any) -> None:
        nonlocal component_counter
        component_counter += 1
        scope = section_stack[-1].id if section_stack else "document"
        path_parts = section_stack[-1].path if section_stack else []
        doc.components.append(
            Component(
                id=stable_id("gt", component_counter),
                type=component_type,
                order=next_order(),
                scope_id=kwargs.pop("scope_id", scope),
                section_path=list(kwargs.pop("section_path", path_parts)),
                **kwargs,
            )
        )

    for child in _iter_docx_body_blocks(body):
        tag = _local_name(child.tag)
        if tag == "p":
            text = _paragraph_text(child)
            if not text:
                _add_figures_from_paragraph(child, rels, media_hashes, add_component)
                continue
            level = _heading_level(child)
            if level:
                while section_stack and section_stack[-1].level >= level:
                    section_stack.pop()
                section_counter += 1
                parent_id = section_stack[-1].id if section_stack else None
                path_parts = [*section_stack[-1].path, text] if section_stack else [text]
                sec = SectionNode(
                    id=stable_id("gt_sec", section_counter),
                    text=text,
                    level=level,
                    parent_id=parent_id,
                    path=path_parts,
                    order=next_order(),
                    scope_id=stable_id("gt_sec", section_counter),
                )
                doc.sections.append(sec)
                section_stack.append(sec)
                add_component("heading", text=text, meta={"level": level})
            elif CAPTION_RE.match(text):
                add_component("caption", text=text)
            else:
                add_component("text", text=text)
            _add_figures_from_paragraph(child, rels, media_hashes, add_component)
        elif tag == "tbl":
            table_html = _docx_table_to_html(child)
            add_component("table", html=table_html, text=_strip_tags(table_html))
    return doc


def _iter_docx_body_blocks(parent: ET.Element) -> Iterable[ET.Element]:
    """Yield paragraph/table blocks in document order, including SDT content.

    Tables in Word files are often wrapped in structured document tags (`w:sdt`).
    We recurse into those containers but do not recurse through `w:tbl` itself,
    because table paragraphs are scored by the table metric, not text fidelity.
    """

    for child in list(parent):
        tag = _local_name(child.tag)
        if tag in {"p", "tbl"}:
            yield child
        elif tag in {"sdt", "sdtContent", "ins"}:
            yield from _iter_docx_body_blocks(child)


def _heading_level(p_elem: ET.Element) -> int:
    ppr = p_elem.find(W + "pPr")
    if ppr is None:
        return 0
    style = ppr.find(W + "pStyle")
    if style is None:
        return 0
    val = style.attrib.get(W + "val") or style.attrib.get("val") or ""
    m = re.search(r"heading\s*(\d+)|Heading(\d+)", val, re.IGNORECASE)
    if m:
        return int(m.group(1) or m.group(2))
    return 0


def _paragraph_text(p_elem: ET.Element) -> str:
    parts = [t.text or "" for t in p_elem.iter(W + "t")]
    return normalize_text("".join(parts))


def _docx_table_to_html(tbl: ET.Element) -> str:
    rows_out: List[str] = ["<table>"]
    active_vmerge: Dict[int, Tuple[str, int]] = {}
    for tr in tbl.findall(W + "tr"):
        rows_out.append("<tr>")
        col = 0
        for tc in tr.findall(W + "tc"):
            text = normalize_text(" ".join(t.text or "" for t in tc.iter(W + "t")))
            tcpr = tc.find(W + "tcPr")
            colspan = 1
            rowspan = 1
            if tcpr is not None:
                grid_span = tcpr.find(W + "gridSpan")
                if grid_span is not None:
                    colspan = _safe_int(grid_span.attrib.get(W + "val") or grid_span.attrib.get("val"), 1)
                vmerge = tcpr.find(W + "vMerge")
                if vmerge is not None:
                    val = vmerge.attrib.get(W + "val") or vmerge.attrib.get("val") or "continue"
                    if val == "continue":
                        # Continuation cells are represented by the originating cell in Word.
                        col += colspan
                        continue
                    active_vmerge[col] = (text, colspan)
            attr = ""
            if colspan > 1:
                attr += f' colspan="{colspan}"'
            if rowspan > 1:
                attr += f' rowspan="{rowspan}"'
            rows_out.append(f"<td{attr}>{html.escape(text)}</td>")
            col += colspan
        rows_out.append("</tr>")
    rows_out.append("</table>")
    return "".join(rows_out)


def _add_figures_from_paragraph(
    p_elem: ET.Element,
    rels: Dict[str, str],
    media_hashes: Dict[str, str],
    add_component,
) -> None:
    for blip in p_elem.iter(A + "blip"):
        rid = blip.attrib.get(R + "embed") or blip.attrib.get(R + "link") or ""
        target = rels.get(rid, "")
        media_name = target.split("/")[-1] if target else ""
        add_component(
            "figure",
            media_path=target,
            relation_id=rid,
            meta={"sha256": media_hashes.get(media_name, "")},
        )


def build_xlsx_reference(path: Path, doc_id: str) -> DocumentComponents:
    try:
        import openpyxl
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("openpyxl is required for XLSX reference extraction") from exc

    wb = openpyxl.load_workbook(path, data_only=False, read_only=False)
    doc = DocumentComponents(
        doc_id=doc_id,
        modality="xlsx",
        source_path=str(path),
        reference_type="native_xlsx_xml",
        gt_strength="strong",
    )
    order = 0
    comp_idx = 0

    def next_order() -> int:
        nonlocal order
        order += 1
        return order

    for sheet_idx, ws in enumerate(wb.worksheets, 1):
        sec = SectionNode(
            id=stable_id("gt_sheet", sheet_idx),
            text=ws.title,
            level=1,
            parent_id=None,
            path=[ws.title],
            order=next_order(),
            scope_id=stable_id("gt_sheet", sheet_idx),
        )
        doc.sections.append(sec)
        for region_idx, (row_start, col_start, rows) in enumerate(_xlsx_non_empty_regions(ws), 1):
            comp_idx += 1
            html_value = _xlsx_rows_to_html(rows, ws.merged_cells.ranges, row_start=row_start, col_start=col_start)
            doc.components.append(
                Component(
                    id=stable_id("gt", comp_idx),
                    type="table",
                    text=_strip_tags(html_value),
                    html=html_value,
                    order=next_order(),
                    scope_id=sec.id,
                    section_path=sec.path,
                )
            )
    return doc


def _xlsx_non_empty_regions(ws: Any) -> List[Tuple[int, int, List[List[str]]]]:
    non_empty = set()
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is not None and normalize_text(str(cell.value)):
                non_empty.add((cell.row, cell.column))
    regions: List[Tuple[int, int, List[List[str]]]] = []
    seen = set()
    for start in sorted(non_empty):
        if start in seen:
            continue
        stack = [start]
        seen.add(start)
        cells = []
        while stack:
            r, c = stack.pop()
            cells.append((r, c))
            for nxt in ((r - 1, c), (r + 1, c), (r, c - 1), (r, c + 1)):
                if nxt in non_empty and nxt not in seen:
                    seen.add(nxt)
                    stack.append(nxt)
        if len(cells) < 2:
            continue
        min_r = min(r for r, _ in cells)
        max_r = max(r for r, _ in cells)
        min_c = min(c for _, c in cells)
        max_c = max(c for _, c in cells)
        rows: List[List[str]] = []
        for r in range(min_r, max_r + 1):
            rows.append([normalize_text(str(ws.cell(r, c).value or "")) for c in range(min_c, max_c + 1)])
        regions.append((min_r, min_c, rows))
    return regions


def _xlsx_rows_to_html(rows: Sequence[Sequence[str]], merge_ranges: Iterable[Any], row_start: int = 1, col_start: int = 1) -> str:
    merge_lookup: Dict[Tuple[int, int], Tuple[int, int]] = {}
    skip: set[Tuple[int, int]] = set()
    for rng in merge_ranges:
        min_col, min_row, max_col, max_row = rng.bounds
        local_start = (min_row - row_start + 1, min_col - col_start + 1)
        merge_lookup[local_start] = (max_row - min_row + 1, max_col - min_col + 1)
        for r in range(min_row, max_row + 1):
            for c in range(min_col, max_col + 1):
                if (r, c) != (min_row, min_col):
                    skip.add((r - row_start + 1, c - col_start + 1))
    out = ["<table>"]
    for r_idx, row in enumerate(rows, 1):
        out.append("<tr>")
        for c_idx, value in enumerate(row, 1):
            if (r_idx, c_idx) in skip:
                continue
            rowspan, colspan = merge_lookup.get((r_idx, c_idx), (1, 1))
            attr = ""
            if rowspan > 1:
                attr += f' rowspan="{rowspan}"'
            if colspan > 1:
                attr += f' colspan="{colspan}"'
            out.append(f"<td{attr}>{html.escape(value)}</td>")
        out.append("</tr>")
    out.append("</table>")
    return "".join(out)


def build_pptx_reference(path: Path, doc_id: str) -> DocumentComponents:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is required for PPTX reference extraction") from exc

    prs = Presentation(path)
    doc = DocumentComponents(
        doc_id=doc_id,
        modality="pptx",
        source_path=str(path),
        reference_type="native_pptx_xml",
        gt_strength="strong",
    )
    order = 0
    comp_idx = 0

    def next_order() -> int:
        nonlocal order
        order += 1
        return order

    for slide_idx, slide in enumerate(prs.slides, 1):
        sec = SectionNode(
            id=stable_id("gt_slide", slide_idx),
            text=f"Slide {slide_idx}",
            level=1,
            parent_id=None,
            path=[f"Slide {slide_idx}"],
            order=next_order(),
            scope_id=stable_id("gt_slide", slide_idx),
        )
        doc.sections.append(sec)
        shapes = sorted(slide.shapes, key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)))
        for shape in shapes:
            comp_idx += 1
            if getattr(shape, "has_table", False):
                html_value = _pptx_table_to_html(shape.table)
                doc.components.append(Component(stable_id("gt", comp_idx), "table", text=_strip_tags(html_value), html=html_value, order=next_order(), scope_id=sec.id, section_path=sec.path))
            elif getattr(shape, "has_text_frame", False):
                text = normalize_text(getattr(shape, "text", "") or "")
                if not text:
                    continue
                ctype = "caption" if CAPTION_RE.match(text) else "text"
                doc.components.append(Component(stable_id("gt", comp_idx), ctype, text=text, order=next_order(), scope_id=sec.id, section_path=sec.path))
            elif getattr(shape, "shape_type", None) is not None and "PICTURE" in str(getattr(shape, "shape_type", "")):
                doc.components.append(Component(stable_id("gt", comp_idx), "figure", order=next_order(), scope_id=sec.id, section_path=sec.path))
    return doc


def _pptx_table_to_html(table: Any) -> str:
    out = ["<table>"]
    for row in table.rows:
        out.append("<tr>")
        for cell in row.cells:
            out.append(f"<td>{html.escape(normalize_text(cell.text or ''))}</td>")
        out.append("</tr>")
    out.append("</table>")
    return "".join(out)


def load_omnidocbench_reference(gt_path: Path, doc_id: Optional[str] = None) -> List[DocumentComponents]:
    payload = json.loads(gt_path.read_text(encoding="utf-8"))
    pages = payload if isinstance(payload, list) else payload.get("pages") or payload.get("data") or []
    docs_by_id: Dict[str, DocumentComponents] = {}
    counters: Dict[str, int] = {}
    for page in pages:
        page_info = page.get("page_info") or {}
        image_path = str(page_info.get("image_path") or "")
        current_doc_id = Path(image_path).stem.split("_")[0] if image_path else str(page_info.get("doc_id") or "omnidocbench")
        if doc_id and current_doc_id != doc_id:
            continue
        doc = docs_by_id.setdefault(
            current_doc_id,
            DocumentComponents(
                doc_id=current_doc_id,
                modality="pdf",
                source_path=str(gt_path),
                reference_type="omnidocbench_annotation",
                gt_strength="strong",
            ),
        )
        scope = f"page_{page_info.get('page_no', len(doc.sections) + 1)}"
        if not any(s.id == scope for s in doc.sections):
            doc.sections.append(SectionNode(scope, scope, 1, None, [scope], len(doc.sections) + 1, scope))
        for det in sorted(page.get("layout_dets") or [], key=lambda d: d.get("order", 10**9)):
            if det.get("ignore"):
                continue
            category = str(det.get("category_type") or "")
            comp_type = _omni_category_to_type(category)
            if comp_type is None:
                continue
            counters[current_doc_id] = counters.get(current_doc_id, 0) + 1
            text = normalize_text(str(det.get("text") or ""))
            html_value = str(det.get("html") or "")
            doc.components.append(
                Component(
                    id=stable_id("gt", counters[current_doc_id]),
                    type=comp_type,
                    text=text or _strip_tags(html_value),
                    html=html_value,
                    order=int(det.get("order", counters[current_doc_id])),
                    scope_id=scope,
                    section_path=[scope],
                    bbox=det.get("poly"),
                    meta={"category_type": category, "anno_id": det.get("anno_id")},
                )
            )
    return list(docs_by_id.values())


def _omni_category_to_type(category: str) -> Optional[str]:
    c = category.lower()
    if "table_caption" in c or "figure_caption" in c:
        return "caption"
    if c == "table":
        return "table"
    if c in {"figure", "image"} or "figure" == c:
        return "figure"
    if any(key in c for key in ("text", "title", "header", "footer", "reference", "equation")):
        return "heading" if "title" in c else "text"
    return None


def _read_rels(zf: zipfile.ZipFile, rels_path: str) -> Dict[str, str]:
    try:
        root = ET.fromstring(zf.read(rels_path))
    except KeyError:
        return {}
    out = {}
    for rel in root:
        rid = rel.attrib.get("Id", "")
        target = rel.attrib.get("Target", "")
        if rid:
            out[rid] = target
    return out


def _media_hashes(zf: zipfile.ZipFile, prefix: str) -> Dict[str, str]:
    out = {}
    for name in zf.namelist():
        if name.startswith(prefix):
            out[Path(name).name] = hashlib.sha256(zf.read(name)).hexdigest()
    return out


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _safe_int(value: Any, default: int = 1) -> int:
    try:
        return int(value)
    except Exception:
        return default


def _strip_tags(value: str) -> str:
    return normalize_text(re.sub(r"<[^>]+>", " ", value or ""))
