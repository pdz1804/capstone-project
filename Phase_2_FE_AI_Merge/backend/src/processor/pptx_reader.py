"""PPTX reader/parser.

Extract structured content (text, images, tables, shapes, and diagrams) from
PowerPoint (.pptx) files and convert it to a hierarchical tree suitable for
JSON output. Key features:
- Robust text extraction (filters strikethrough runs).
- TOC (Table of Contents) detection and optional skipping (supports English/Japanese).
- Image and diagram capture (extracts image blobs or crops exported slide images).
- Table extraction with markdown formatting and multi-slide continuation merging.
- Diagram grouping via spatial clustering and connector-aware merging.
- Header detection and optional LLM-based validation; builds a ContentNode tree.
- Main classes/functions: PptxParser, extract_pptx_text (convenience wrapper).

Dependencies: python-pptx, Pillow, numpy, scikit-learn. Optional: libreoffice, pdf2image/poppler for slide export.
"""
import os
import re
import io
import hashlib
import tempfile
import shutil
import subprocess
import json
import importlib.util
from dataclasses import dataclass, field
from typing import List, Dict, Optional, Tuple, Any
from enum import Enum
from pathlib import Path
import requests
import sys

_THIS_DIR = Path(__file__).resolve().parent
_UTILS_DIR = _THIS_DIR / "utils"


def _load_local_module(module_name: str, file_name: str):
    spec = importlib.util.spec_from_file_location(module_name, _UTILS_DIR / file_name)
    if spec is None or spec.loader is None:
        raise ImportError(f"Could not load local module {file_name}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


try:
    from .utils.heading_detector import detect_heading_level, normalize_line
except Exception:
    _heading_detector = _load_local_module("pptx_reader_heading_detector", "heading_detector.py")
    detect_heading_level = _heading_detector.detect_heading_level
    normalize_line = _heading_detector.normalize_line

try:
    from .utils.phash import phash, similarity
except Exception:
    _phash = _load_local_module("pptx_reader_phash", "phash.py")
    phash = _phash.phash
    similarity = _phash.similarity

try:
    from .utils.gen_s3_url import get_signed_url  # type: ignore
except Exception:
    def get_signed_url(path: str) -> str:
        return path

try:
    from .utils.ignored_image_router import get_image_urls  # type: ignore
except Exception:
    def get_image_urls(repository_id: str = "") -> List[str]:
        return []

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.connector import Connector
from pptx.table import Table
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn 
from PIL import Image
import cv2
import numpy as np
from sklearn.cluster import DBSCAN
import time
# from app.cm_shared.llm_provider import LLMProvider

# THIS BLOCK IS FOR LOCAL TESTING PURPOSES ONLY.

# Try normal import first; if it fails, attempt to make the project root importable and retry.
try:
    from app.cm_shared.llm_provider import LLMProvider
except Exception as _import_err:
    import sys, os
    # Attempt to locate project root (the directory that contains 'app')
    try:
        this_file = Path(__file__).resolve()
    except Exception:
        this_file = None

    project_root = None
    if this_file:
        cur = this_file
        for _ in range(8):
            cur = cur.parent
            if (cur / "app").is_dir():
                project_root = cur
                break

    # fallback to env var if set
    if project_root is None:
        env_root = os.getenv("PROJECT_ROOT")
        if env_root and Path(env_root).is_dir():
            project_root = Path(env_root)

    if project_root:
        sys.path.insert(0, str(project_root))

    try:
        from app.cm_shared.llm_provider import LLMProvider
    except Exception as e:
        # Provide a clear runtime error when LLMProvider is actually used but unavailable
        class _MissingLLMProvider:
            @staticmethod
            def create(*args, **kwargs):
                raise RuntimeError(
                    "LLMProvider could not be imported. Original error: "
                    + repr(_import_err) + " ; Retry error: " + repr(e)
                )
        LLMProvider = _MissingLLMProvider

IMAGE_PATH_START_MARKER = "[START_IMAGE_PATH]"
IMAGE_PATH_END_MARKER = "[END_IMAGE_PATH]"
TABLE_CONTENT_START_MARKER = "[START_TABLE_CONTENT]"
TABLE_CONTENT_END_MARKER = "[END_TABLE_CONTENT]"
SHAPE_CONTENT_START_MARKER = "[START_SHAPE_CONTENT]"
SHAPE_CONTENT_END_MARKER = "[END_SHAPE_CONTENT]"
IMAGE_CONTENT_START_MARKER = "[START_IMAGE_CONTENT]"
IMAGE_CONTENT_END_MARKER = "[END_IMAGE_CONTENT]"

class ElementType(Enum):
    """Types of elements in a PowerPoint slide."""
    TEXTBOX = "textbox"
    IMAGE = "image"
    TABLE = "table"
    SHAPE = "shape"  # Geometric shapes (diamond, rectangle, etc.)
    GROUP = "group"  # Grouped elements (diagrams)
    CONNECTOR = "connector"
    PLACEHOLDER = "placeholder"
    UNKNOWN = "unknown"


@dataclass
class BoundingBox:
    """
    Represents a bounding box with position and size.
    BBox in EMUs
    """
    x: int   
    y: int   
    width: int  
    height: int  

    @property
    def x2(self) -> int:
        return self.x + self.width

    @property
    def y2(self) -> int:
        return self.y + self.height

    @property
    def center_x(self) -> int:
        return self.x + self.width // 2

    @property
    def center_y(self) -> int:
        return self.y + self.height // 2

    def overlaps(self, other: "BoundingBox", threshold: float = 0.0) -> bool:
        """Check if two bounding boxes overlap."""
        if threshold > 0:
            x_overlap = max(0, min(self.x2, other.x2) - max(self.x, other.x))
            y_overlap = max(0, min(self.y2, other.y2) - max(self.y, other.y))
            intersection = x_overlap * y_overlap
            min_area = min(self.width * self.height, other.width * other.height)
            return (intersection / min_area) >= threshold if min_area > 0 else False
        return not (self.x2 < other.x or other.x2 < self.x or
                    self.y2 < other.y or other.y2 < self.y)

    def get_overlap_ratio(self, other: "BoundingBox") -> float:
        """
        Calculate the overlap ratio (intersection area / this bbox area).
        
        Returns:
            Float from 0.0 to 1.0 representing how much of this box is covered by other.
        """
        x_overlap = max(0, min(self.x2, other.x2) - max(self.x, other.x))
        y_overlap = max(0, min(self.y2, other.y2) - max(self.y, other.y))
        intersection = x_overlap * y_overlap
        self_area = self.width * self.height
        return intersection / self_area if self_area > 0 else 0.0


    def contains_point(self, x: int, y: int) -> bool:
        """Check if a point is inside this bounding box."""
        return self.x <= x <= self.x2 and self.y <= y <= self.y2

    def union(self, other: "BoundingBox") -> "BoundingBox":
        """Return the union bounding box of two boxes."""
        new_x = min(self.x, other.x)
        new_y = min(self.y, other.y)
        new_x2 = max(self.x2, other.x2)
        new_y2 = max(self.y2, other.y2)
        return BoundingBox(new_x, new_y, new_x2 - new_x, new_y2 - new_y)

    @classmethod
    def from_shape(cls, shape: BaseShape) -> "BoundingBox":
        """Create BoundingBox from a shape."""
        return cls(
            x=shape.left or 0,
            y=shape.top or 0,
            width=shape.width or 0,
            height=shape.height or 0
        )


@dataclass
class LineBBox:
    """Represents a line of text with its bounding box."""
    text: str
    bbox: BoundingBox
    
    def is_covered_by(self, overlay_bbox: BoundingBox, vertical_threshold: float = 0.5, horizontal_threshold: float = 0.3) -> bool:
        """
        Check if this line is covered by an overlay bounding box.
        
        Args:
            overlay_bbox: The bounding box of the overlaying shape
            vertical_threshold: Minimum vertical overlap ratio (0-1) to consider line covered
            horizontal_threshold: Minimum horizontal overlap ratio (0-1) to consider line covered
            
        Returns:
            True if line is sufficiently covered by the overlay
        """
        # Calculate vertical overlap
        v_overlap = max(0, min(self.bbox.y2, overlay_bbox.y2) - max(self.bbox.y, overlay_bbox.y))
        v_ratio = v_overlap / self.bbox.height if self.bbox.height > 0 else 0
        
        # Calculate horizontal overlap
        h_overlap = max(0, min(self.bbox.x2, overlay_bbox.x2) - max(self.bbox.x, overlay_bbox.x))
        h_ratio = h_overlap / self.bbox.width if self.bbox.width > 0 else 0
        
        return v_ratio >= vertical_threshold and h_ratio >= horizontal_threshold


@dataclass
class SlideElement:
    """Represents an extracted element from a slide."""
    element_type: ElementType
    bbox: BoundingBox
    content: str = ""
    raw_shape: Optional[BaseShape] = None
    children: List["SlideElement"] = field(default_factory=list)
    shape_type_name: str = ""  
    image_path: Optional[str] = None
    image_hash: Optional[str] = None
    table_data: Optional[List[List[str]]] = None
    is_merged: bool = False
    slide_index: int = 0
    z_index: int = 0  
    # is_deprecated: bool = False
    is_header: bool = False
    # Line-level bounding boxes for accurate deprecation detection
    line_bboxes: Optional[List[LineBBox]] = None
    # Table continuation tracking fields
    table_column_count: Optional[int] = None
    table_column_widths: Optional[List[float]] = None  # Normalized widths (0-1)
    table_first_row: Optional[List[str]] = None  # First row for header matching
    table_id: Optional[str] = None  # Unique identifier for this table
    is_table_continuation: bool = False  # True if this is a continuation of a previous table
    continuation_of_table_id: Optional[str] = None  # ID of the parent table
    merged_from_slides: Optional[List[int]] = None  # List of slide indices this table spans

    def __lt__(self, other: "SlideElement") -> bool:
        """Sort by y then x position."""
        if self.bbox.y != other.bbox.y:
            return self.bbox.y < other.bbox.y
        return self.bbox.x < other.bbox.x


@dataclass
class HeaderCandidate:
    """Represents a potential header detected in content."""
    text: str
    pattern_match: str
    element_index: int
    position_score: float  # low y ==> high score
    is_validated: bool = False


@dataclass
class ContentNode:
    """Represents a node in the output tree structure."""
    heading_text: str
    heading_level: int
    content: str
    children: List["ContentNode"] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary format."""
        return {
            "heading_text": self.heading_text,
            "heading_level": self.heading_level,
            "content": self.content,
            "children": [child.to_dict() for child in self.children]
        }


class PptxParser:
    """
    Parser for extracting structured content from PowerPoint files.

    Extracts text, images, tables, and diagrams from PPTX files and
    converts them to a hierarchical JSON structure with detected headers.
    """

    # -------------------- TOC Detection Configuration --------------------
    # Common TOC header patterns for English and Japanese
    _TOC_HEADER_PATTERNS = [
        # Exact matches (strict)
        re.compile(r'^\s*table\s+of\s+contents\s*$', re.IGNORECASE),
        re.compile(r'^\s*contents\s*$', re.IGNORECASE),
        re.compile(r'^\s*index\s*$', re.IGNORECASE),
        re.compile(r'^\s*[【\[\(（]?\s*目\s*次\s*[】\]\)）]?\s*$'),
        # Flexible matches - detect "目次" or "Table of Contents" anywhere in line
        # This handles numbered TOC headers like "0 目次(Table of Contents)"
        re.compile(r'^\s*\d*\s*目\s*次.*$'),  # "0 目次..." or "目次..."
        re.compile(r'.*table\s+of\s+contents.*', re.IGNORECASE),  # Contains "Table of Contents"
        re.compile(r'^\s*[【\[\(（]?\s*索\s*引\s*[】\]\)）]?\s*$'),
        re.compile(r'^\s*[【\[\(（]?\s*もくじ\s*[】\]\)）]?\s*$'),
        re.compile(r'^\s*[【\[\(（]?\s*モクジ\s*[】\]\)）]?\s*$'),
        re.compile(r'^\s*CONTENTS\s*$', re.IGNORECASE),
    ]

    # TOC entry patterns - lines that look like TOC entries
    _TOC_ENTRY_PATTERNS = [
        # Pattern: "Title ........... 123" or "Title ... 123" (dot leaders with spaces between dots)
        re.compile(r'(?=.*[A-Za-z\u3040-\u30FF\u4E00-\u9FFF]).+[.\s]{5,}\d+\s*$'),
        # Pattern: "1.2.3 Title ... 45" (numbered with dots leading to page)
        re.compile(r'^[\d.]+\s+.+[.\s]{5,}\d+\s*$'),
        # Pattern with explicit dot sequences (. . . . or .......)
        re.compile(r'.+(?:\.[\s.]*){4,}\s*\d+\s*$'),
        # Pattern: "Title  123" at end of line with multiple spaces before page number
        re.compile(r'^(?=.{10,})(?=.*[A-Za-z\u3040-\u30FF\u4E00-\u9FFF]).*\s{3,}\d{1,4}\s*$'),
        # Japanese style with full-width dots and numbers
        re.compile(r'(?=.*[\u3040-\u30FF\u4E00-\u9FFF]).+[・．…\s.]{3,}\s*[\d０-９]+\s*$'),
    ]

    # Patterns that indicate this is likely a TOC entry (not a real heading)
    _TOC_LINE_PATTERNS = [
        re.compile(r'\.{3,}\s*\d+\s*$'),  # Dot leaders followed by page number
        re.compile(r'\s{3,}\d+\s*$'),  # Multiple spaces followed by page number
        re.compile(r'[・．…]{3,}\s*\d+\s*$'),  # Japanese dot leaders
    ]

    def __init__(
        self,
        get_table_content: bool = True,
        get_image_content: bool = True,
        get_shape_content: bool = True,
        validate_headers_with_llm: bool = True,
        merge_threshold_width_ratio: float = 0.3,  
        merge_threshold_content_length: int = 50,  
        keep_temp_files: bool = False,
        merge_continuation_tables: bool = True,
        table_similarity_threshold: float = 0.75,
        table_column_width_tolerance: float = 0.15,
        skip_toc: bool = True,
        repository_id: str = ""
    ):
        """
        Initialize the PPTX parser.
        Args:
            get_table_content: Whether to extract full table content
            get_image_content: Whether to extract images
            get_shape_content: Whether to extract shape content
            validate_headers_with_llm: Whether to use LLM for header validation
            merge_threshold_width_ratio: Width ratio threshold for merge candidates
            merge_threshold_content_length: Content length threshold for merge candidates
            keep_temp_files: Whether to keep temp files (images) after processing
            merge_continuation_tables: Whether to merge tables spanning multiple slides
            table_similarity_threshold: Threshold (0-1) for detecting table continuations
            table_column_width_tolerance: Tolerance for column width matching (0-1)
            skip_toc: Whether to skip Table of Contents slides during extraction
        """
        self.get_table_content = get_table_content
        self.get_image_content = get_image_content
        self.get_shape_content = get_shape_content
        self.validate_headers_with_llm = validate_headers_with_llm
        self.merge_threshold_width_ratio = merge_threshold_width_ratio
        self.merge_threshold_content_length = merge_threshold_content_length
        self.keep_temp_files = keep_temp_files
        self.merge_continuation_tables = merge_continuation_tables
        self.table_similarity_threshold = table_similarity_threshold
        self.table_column_width_tolerance = table_column_width_tolerance
        self._repository_id = repository_id
        self.ignored_images = self._get_ignored_images()
        self.skip_toc = skip_toc
        self._table_id_counter: int = 0

        self._temp_dir: Optional[str] = None
        self._image_counter: int = 0
        self._saved_images: List[str] = []
        self._llm_client = None
        self._slide_width: int = 0
        self._slide_height: int = 0
        self._pptx_path: Optional[str] = None
        self._slide_exports_dir: Optional[str] = None
        self._slide_export_cache: Dict[int, str] = {}
        self._toc_slides: set = set()
        self._marker_freq_map: Dict[str, int] = {}
        self._marker_scan_count: int = 0
        self._extracted_element: List[SlideElement] = []


        
        
        

    # -------------------- TOC Detection Methods --------------------
    
    def _is_toc_header_line(self, text: str) -> bool:
        """Check if text matches a TOC header pattern."""
        text = text.strip()
        for pattern in self._TOC_HEADER_PATTERNS:
            if pattern.match(text):
                return True
        return False

    def _is_toc_entry_line(self, text: str) -> bool:
        """
        Check if a line looks like a TOC entry (title with page number).
        """
        text = text.strip()
        if not text or len(text) < 5:
            return False
        
        for pattern in self._TOC_ENTRY_PATTERNS:
            if pattern.match(text):
                return True
        
        # Additional heuristic: check for dot leaders pattern
        # Count dots (including spaced dots like ". . . ." or ".......")
        dot_count = text.count('.')
        
        # Check if line ends with page number (1-4 digits)
        page_num_match = re.search(r'\d{1,4}\s*$', text)
        
        if dot_count >= 4 and page_num_match:
            # Check that there's some title content (not just dots and numbers)
            content = re.sub(r'[.\s\d]+$', '', text).strip()
            if len(content) >= 3:
                return True
        
        # Check for repeated dot-space pattern (. . . . .)
        dot_space_pattern = re.findall(r'\.[\s.]+', text)
        if len(dot_space_pattern) >= 3 and page_num_match:
            return True
        
        return False

    def _is_toc_line(self, text: str) -> bool:
        """Check if a line looks like a TOC entry based on patterns."""
        text = text.strip()
        for pattern in self._TOC_LINE_PATTERNS:
            if pattern.search(text):
                return True
        return False
    
    def _header_distance(self, actual_level: int, detected_level: int)-> int:
        distance = detected_level - actual_level
        return distance

    def _parse_pptx_toc_entry(
        self, text: str, latest_toc_entry: Optional[Dict[str, Any]], is_first: bool, distance_from_origin: int
    ) -> Tuple[Optional[Dict[str, Any]], bool, int]:
        """
        Parse a single PPTX TOC line into a heading entry.

        Unlike PDF TOC entries, PPTX TOC entries typically lack page numbers.
        We extract heading_text and heading_level only.

        Heading level is determined by identifying a shared numbering prefix
        across all entries (e.g., '4-1' in '4-1.2.3') and counting segments
        after the prefix:
        - '4-1.1' with prefix '4-1' → 1 segment after prefix → Level 1
        - '4-1.2.1' with prefix '4-1' → 2 segments after prefix → Level 2
        - '4-1.2.10' with prefix '4-1' → 2 segments after prefix → Level 2

        Returns:
            Dict with keys: heading_text, heading_level
            or None if parsing fails or if line is a TOC header.
        """
        text = text.strip()

        if not text:
            return None, is_first, distance_from_origin

        # Reject TOC header lines
        if self._is_toc_header_line(text):
            print("Reject TOC header line:", text)
            return None, is_first, distance_from_origin

        # Remove page number and dot leaders / separators
        # This regex matches dots, spaces, etc. at the end followed by the digits
        heading_text = re.sub(r'[.．…・]{2,}\s*\d*\s*$', '', text).strip()
        if not heading_text:
            return None, is_first, distance_from_origin

        # Reject pure annotation lines (e.g., "(24MM：CV#2)", "▲j（24MM:1A）")
        if re.match(r'^[▲△▽▼]?\s*[a-zA-Z0-9]?\s*[（(][^）)]+[）)]$', heading_text):
            return None, is_first, distance_from_origin
        # Reject lines that are purely markers like "▲a", "▲j"
        if re.match(r'^[▲△▽▼][a-zA-Z0-9]?$', heading_text):
            return None, is_first, distance_from_origin

        # Reject very short/trivial lines
        if len(heading_text) < 2:
            return None, is_first, distance_from_origin

        # Reject lines that look like footers (e.g., "Copyright TOYOTA MOTOR CORPORATION")
        if re.match(r'^Copyright\s', heading_text, re.IGNORECASE):
            return None, is_first, distance_from_origin

        # Reject number-only entries (e.g., "4-1.2.18" from strikethrough TOC lines
        # where the title text was removed but the number remained)
        cleaned_for_numcheck = heading_text.replace('\t', ' ').strip()
        if re.match(r'^\d+(?:[-.．…・]\d+)*\s*$', cleaned_for_numcheck):
            return None, is_first, distance_from_origin

        # --- Detect level from numbering structure ---
        # Normalize dashes (fullwidth → ASCII for counting)
        normalized = heading_text.replace('－', '-').replace(' ', '-').replace('‐', '-')
        # Also normalize tabs to spaces
        normalized = normalized.replace('\t', ' ')
        level = None
        is_pre_chapter = False

        # If no pattern matches, try the heading_detector as fallback
        normalized_heading_text = normalize_line(heading_text)
         
        detected_level, pattern_name, match = detect_heading_level(normalized_heading_text)
        if detected_level is not None and not is_first:
            level = detected_level - distance_from_origin if detected_level != 1 else detected_level
            if normalized_heading_text[0] == '0':
                is_pre_chapter = True
        else:
            if latest_toc_entry is not None :
                lastest_level = latest_toc_entry["heading_level"]
                if lastest_level == 3:
                    level = 3
                else:
                    level = lastest_level + 1
            else:
                level = 1
                is_first = False
                if detected_level:
                    distance_from_origin = self._header_distance(level, detected_level)
                     
        # Cap at MAX_DEPTH = 3
        if level > 3:
            return None, is_first, distance_from_origin
        result = {
            "heading_text": heading_text,
            "heading_level": level,
            "is_pre_chapter": is_pre_chapter,
        }
        return (result, is_first, distance_from_origin)

    def _extract_text_from_slide(self, slide) -> List[str]:
        """
        Extract all text content from a slide as a list of text lines.
        
        Args:
            slide: A slide object from python-pptx
            
        Returns:
            List of text strings from all shapes in the slide
        """
        text_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = self._extract_text_from_text_frame(shape.text_frame)
                if text:
                    # Split by newlines to get individual lines
                    for line in text.split('\n'):
                        line = line.strip()
                        if line:
                            text_lines.append(line)
            elif hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    for line in text.split('\n'):
                        line = line.strip()
                        if line:
                            text_lines.append(line)
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if hasattr(cell, "text_frame"):
                            text = self._extract_text_from_text_frame(cell.text_frame)
                        else:
                            text = cell.text.strip() if hasattr(cell, "text") else ""
                        if text:
                            for line in text.split('\n'):
                                line = line.strip()
                                if line:
                                    text_lines.append(line)
        return text_lines

    def _detect_toc_slides(self, prs: Presentation) -> Tuple[set, List[Dict[str, Any]]]:
        """
        Detect which slides are Table of Contents slides and extract TOC entries.
        
        Strategy:
        1. Find slides with TOC header ("目次", "Table of Contents", etc.)
        2. Mark that slide and consecutive slides with TOC-like entries as TOC slides
        3. Stop when we encounter a slide without TOC patterns
        4. Extract heading_text + heading_level from each TOC entry line
        
        Args:
            prs: Presentation object from python-pptx
            
        Returns:
            Tuple of (toc_slide_indices, toc_entries)
            where toc_entries is a list of dicts with heading_text, heading_level
        """
        toc_slides = set()
        toc_entries: List[Dict[str, Any]] = []
        is_first = True
        distance_from_origin = 0
        if not self.skip_toc:
            return toc_slides, toc_entries
        
        def _collect_entries_from_lines(text_lines: List[str], is_first: bool, distance_from_origin: int):
            """Extract TOC entries from text lines on a TOC slide."""
            
            for line in text_lines:
                line = line.strip()
                if not line or len(line) < 3:
                    continue
                # Skip the TOC header itself
                if self._is_toc_header_line(line):
                    continue
                # Parse as TOC entry (works with or without page numbers/dot leaders)
                entry, first_value, distance = self._parse_pptx_toc_entry(
                    line,
                    toc_entries[-1] if toc_entries else None,
                    is_first,
                    distance_from_origin
                )
                is_first = first_value
                distance_from_origin = distance
                if entry:
                    toc_entries.append(entry)
        
        try:
            in_toc_section = False
            toc_entry_threshold = 0.25  # At least 25% of lines should look like TOC entries
            
            for slide_idx, slide in enumerate(prs.slides):
                # Extract text lines from the slide
                slide_text_lines = self._extract_text_from_slide(slide)
                
                if not slide_text_lines:
                    if in_toc_section:
                        # Empty slide in TOC section might be intentional break
                        in_toc_section = False
                    continue
                
                # Check for TOC header on this slide
                has_toc_header = False
                for line in slide_text_lines[:10]:  # Check first 10 lines
                    if self._is_toc_header_line(line):
                        has_toc_header = True
                        in_toc_section = True
                        print(f"[TOC_DEBUG] Found TOC header on slide {slide_idx + 1}: '{line}'")
                        break
                
                if has_toc_header:
                    toc_slides.add(slide_idx)
                    print(f"[TOC] Detected TOC header on slide {slide_idx + 1}")
                    _collect_entries_from_lines(slide_text_lines, is_first, distance_from_origin)
                    continue
                
                # If we're in a TOC section, check if this slide continues it
                if in_toc_section:
                    # Count how many lines look like TOC entries
                    toc_entry_count = 0
                    total_content_lines = 0
                    
                    for line in slide_text_lines:
                        line = line.strip()
                        if not line or len(line) < 3:
                            continue
                        total_content_lines += 1
                        if self._is_toc_entry_line(line) or self._is_toc_line(line):
                            toc_entry_count += 1
                        else:
                            # Relaxed check: if line matches a heading pattern,
                            # it's likely a PPTX TOC entry without dot leaders
                            lvl, _, _ = detect_heading_level(normalize_line(line))
                            if lvl is not None:
                                toc_entry_count += 1
                    
                    # If significant portion of slide is TOC entries, mark as TOC slide
                    if total_content_lines > 0:
                        toc_ratio = toc_entry_count / total_content_lines
                        if toc_ratio >= toc_entry_threshold:
                            toc_slides.add(slide_idx)
                            print(f"[TOC] Slide {slide_idx + 1} detected as TOC continuation (ratio: {toc_ratio:.2f}, {toc_entry_count}/{total_content_lines})")
                            _collect_entries_from_lines(slide_text_lines, is_first, distance_from_origin)
                        else:
                            # TOC section ended
                            in_toc_section = False
                            print(f"[TOC] TOC section ended at slide {slide_idx + 1} (ratio: {toc_ratio:.2f}, {toc_entry_count}/{total_content_lines})")
                else:
                    # Not in TOC section - check if this is a standalone TOC slide
                    # (some presentations have TOC header on same slide as entries)
                    toc_entry_count = 0
                    total_content_lines = 0
                    
                    for line in slide_text_lines:
                        line = line.strip()
                        if not line or len(line) < 3:
                            continue
                        total_content_lines += 1
                        if self._is_toc_entry_line(line) or self._is_toc_line(line):
                            toc_entry_count += 1
                    
                    # Higher threshold for standalone detection (50%)
                    if total_content_lines > 5 and toc_entry_count / total_content_lines >= 0.5:
                        toc_slides.add(slide_idx)
                        in_toc_section = True
                        print(f"[TOC] Slide {slide_idx + 1} detected as TOC by content pattern (ratio: {toc_entry_count/total_content_lines:.2f})")
                        _collect_entries_from_lines(slide_text_lines, is_first, distance_from_origin)
                        
        except Exception as exc:
            print(f"[TOC_ERROR] Error detecting TOC slides: {exc}")
            import traceback
            traceback.print_exc()
        
        print(f"[TOC] Detected {len(toc_slides)} TOC slides, {len(toc_entries)} TOC entries")
        return toc_slides, toc_entries

    def _is_toc_slide(self, slide_index: int) -> bool:
        """
        Check if a slide (0-based index) is a TOC slide.
        """
        return slide_index in self._toc_slides

    # -------------------- End TOC Detection Methods --------------------

    def _get_temp_dir(self) -> str:
        """Get or create temporary directory for images."""
        if self._temp_dir is None:
            self._temp_dir = tempfile.mkdtemp(prefix="pptx_images_")
        return self._temp_dir

    def _compute_md5(self, data: bytes) -> str:
        """Compute MD5 hash of data."""
        return hashlib.md5(data).hexdigest()

    def _save_image(self, image_bytes: bytes, extension: str = "png") -> Tuple[str, str]:
        """
        Save image to temp directory and return path and hash.

        Returns:
            Tuple of (image_path, md5_hash)
        """
        md5_hash = self._compute_md5(image_bytes)
        self._image_counter += 1
        filename = f"img_{self._image_counter}_{md5_hash[:8]}.{extension}"
        filepath = os.path.join(self._get_temp_dir(), filename)

        with open(filepath, "wb") as f:
            f.write(image_bytes)

        self._saved_images.append(filepath)
        return filepath, md5_hash

    def _get_slide_exports_dir(self) -> str:
        """Get or create temporary directory for slide exports."""
        if self._slide_exports_dir is None:
            self._slide_exports_dir = tempfile.mkdtemp(prefix="pptx_slides_")
        return self._slide_exports_dir

    def _export_slides_as_images(self) -> bool:
        """
        Export all slides as PNG images using LibreOffice (via PDF).

        Converts PPTX to PDF first, then extracts each page as a PNG image.

        Returns:
            True if export succeeded, False otherwise.
        """
        if not self._pptx_path:
            return False

        if self._slide_export_cache:
            return True

        export_dir = self._get_slide_exports_dir()
        pptx_basename = Path(self._pptx_path).stem

        try:
            pdf_path = os.path.join(export_dir, f"{pptx_basename}.pdf")

            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", export_dir,
                    self._pptx_path,
                ],
                capture_output=True,
                timeout=120,
            )

            if result.returncode != 0:
                print(f"Warning: LibreOffice PDF conversion failed: {result.stderr.decode()}")
                return False

            if not os.path.exists(pdf_path):
                print("Warning: PDF file was not created")
                return False

            # Method 1: pdftoppm (PREFERRED - writes directly to disk, no RAM bloat)
            # This is the most memory-efficient approach as it doesn't load images into Python
            try:
                result = subprocess.run(
                    [
                        "pdftoppm",
                        "-png",
                        "-r", "150",  
                        pdf_path,
                        os.path.join(export_dir, "slide"),
                    ],
                    capture_output=True,
                    timeout=300,  # Longer timeout for large presentations
                )

                if result.returncode == 0:
                    for png_file in sorted(Path(export_dir).glob("slide-*.png")):
                        try:
                            page_num = int(png_file.stem.split("-")[-1])
                            self._slide_export_cache[page_num - 1] = str(png_file)
                        except ValueError:
                            pass
                    if self._slide_export_cache:
                        return True
            except FileNotFoundError:
                pass  # pdftoppm not installed, try next method
            except subprocess.TimeoutExpired:
                print("Warning: pdftoppm conversion timed out")

            # Method 2: pdf2image with per-page processing (avoids loading all pages at once)
            try:
                from pdf2image import convert_from_path
                from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError
                
                # Get page count first to process one page at a time
                try:
                    from pdf2image import pdfinfo_from_path
                    info = pdfinfo_from_path(pdf_path)
                    page_count = info.get('Pages', 0)
                except Exception:
                    # Fallback: estimate from file or use a reasonable default
                    page_count = 100  # Will stop when no more pages
                
                # Process ONE page at a time to minimize RAM usage
                for page_num in range(1, page_count + 1):
                    try:
                        # Convert single page only
                        images = convert_from_path(
                            pdf_path, 
                            dpi=150,
                            first_page=page_num,
                            last_page=page_num,
                            thread_count=1  # Single thread to limit memory
                        )
                        if images:
                            img = images[0]
                            png_path = os.path.join(export_dir, f"slide_{page_num}.png")
                            img.save(png_path, "PNG")
                            img.close()  # Explicitly close to free memory
                            del img
                            self._slide_export_cache[page_num - 1] = png_path
                        else:
                            break  # No more pages
                    except Exception as page_err:
                        # If we already have some pages, that's fine
                        if self._slide_export_cache:
                            break
                        raise page_err
                
                if self._slide_export_cache:
                    return True
                    
            except ImportError:
                pass  # pdf2image not installed
            except Exception as e:
                print(f"Warning: pdf2image conversion failed: {e}")

            # Method 3: PIL fallback (least reliable, but doesn't load all at once)
            try:
                with Image.open(pdf_path) as pdf_img:
                    for idx in range(getattr(pdf_img, 'n_frames', 1)):
                        pdf_img.seek(idx)
                        png_path = os.path.join(export_dir, f"slide_{idx + 1}.png")
                        pdf_img.save(png_path, "PNG")
                        self._slide_export_cache[idx] = png_path
                return bool(self._slide_export_cache)
            except Exception:
                pass

            print("Warning: Could not convert PDF to images. Install poppler-utils (pdftoppm) for best results.")
            return False

        except subprocess.TimeoutExpired:
            print("Warning: LibreOffice conversion timed out")
            return False
        except FileNotFoundError:
            print("Warning: LibreOffice not found. Install libreoffice to export diagrams as images.")
            return False
        except Exception as e:
            print(f"Warning: Failed to export slides: {e}")
            return False
    def _get_exported_slide_path(self, slide_index: int) -> Optional[str]:
        """Get the path to an exported slide image."""
        if slide_index in self._slide_export_cache:
            return self._slide_export_cache[slide_index]
        if not self._slide_export_cache:
            self._export_slides_as_images()
        return self._slide_export_cache.get(slide_index)
    def _crop_region_from_slide(
        self,
        slide_image_path: str,
        bbox: BoundingBox
    ) -> Optional[bytes]:
        """
        Crop a region from the exported slide image based on bounding box.

        Args:
            slide_image_path: Path to the exported slide PNG
            bbox: Bounding box in EMU coordinates

        Returns:
            Cropped image bytes or None if failed
        """
        try:
            with Image.open(slide_image_path) as img:
                img_width, img_height = img.size
                
                if self._slide_width <= 0 or self._slide_height <= 0:
                    return None
                
                scale_x = img_width / self._slide_width
                scale_y = img_height / self._slide_height

                print(f"[DEBUG] Original bboxes: left {bbox.x}, top {bbox.y}, right {bbox.x2}, bottom {bbox.y2}")
                left = int(bbox.x * scale_x)
                top = int(bbox.y * scale_y)
                right = int(bbox.x2 * scale_x)
                bottom = int(bbox.y2 * scale_y)
                print(f"[DEBUG] Region bbox identification : left {left}, top {top}, right {right}, bottom {bottom}")
                


                left = max(0, min(left, img_width - 1))
                top = max(0, min(top, img_height - 1))
                right = max(left + 1, min(right, img_width))
                bottom = max(top + 1, min(bottom, img_height))

                cropped = img.crop((left, top, right, bottom))
                
                buffer = io.BytesIO()
                cropped.save(buffer, format="PNG")
                cropped.close()  # Explicitly close cropped image
                 
                return buffer.getvalue()

        except Exception as e:
            print(f"Warning: Failed to crop region from slide: {e}")
            return None

    def _get_element_type(self, shape: BaseShape) -> ElementType:
        """Determine the type of a PowerPoint shape."""
        if isinstance(shape, GroupShape):
            return ElementType.GROUP
        elif isinstance(shape, Picture):
            return ElementType.IMAGE
        elif isinstance(shape, GraphicFrame):
            if shape.has_table:
                return ElementType.TABLE
            return ElementType.UNKNOWN
        elif isinstance(shape, Connector):
            return ElementType.CONNECTOR
        elif isinstance(shape, AutoShape):
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    return ElementType.TEXTBOX
            except Exception:
                pass
            try:
                auto_type = shape.auto_shape_type
                if auto_type is not None:
                    shape_type_name = str(auto_type).split(".")[-1]
                    if "RECTANGLE" in shape_type_name.upper():
                        try:
                            if hasattr(shape, 'fill'):
                                from pptx.enum.dml import MSO_FILL_TYPE
                                if shape.fill.type in (MSO_FILL_TYPE.BACKGROUND, None):
                                    return ElementType.TEXTBOX
                                if shape.fill.type == MSO_FILL_TYPE.SOLID:
                                    try:
                                        rgb = shape.fill.fore_color.rgb
                                        r, g, b = rgb[0], rgb[1], rgb[2]
                                        is_white = r > 240 and g > 240 and b > 240
                                        if is_white:
                                            return ElementType.TEXTBOX
                                        else:
                                            return ElementType.SHAPE
                                    except:
                                        return ElementType.SHAPE
                                return ElementType.SHAPE
                        except:
                            pass
                        return ElementType.TEXTBOX
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        return ElementType.SHAPE  
                    return ElementType.SHAPE
            except (ValueError, AttributeError):
                pass
            if shape.has_text_frame and shape.text_frame.text.strip():
                return ElementType.TEXTBOX
            return ElementType.SHAPE
        try:
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                return ElementType.PLACEHOLDER
        except Exception:
            pass

        return ElementType.UNKNOWN

    def _has_strikethrough(self, run) -> bool:
        """
        Check if a run has strikethrough formatting via XML.
        
        python-pptx doesn't expose strikethrough directly, so we need to
        check the XML element for the 'strike' attribute on <a:rPr>.
        
        Args:
            run: A run object from python-pptx paragraph
            
        Returns:
            True if the run has strikethrough formatting
        """
        try:
            # Access the run's XML element
            r_elem = run._r
            if r_elem is None:
                return False
            
            # Find the <a:rPr> (run properties) element
            rPr = r_elem.find(qn('a:rPr'))
            if rPr is None:
                return False
            
            # Check for 'strike' attribute - values can be:
            # 'sngStrike' (single), 'dblStrike' (double), 'noStrike' (none)
            strike = rPr.get('strike')
            if strike is not None and strike != 'noStrike':
                return True
            
            return False
        except Exception:
            return False

    def _extract_text_from_text_frame(self, text_frame) -> str:
        """
        Extract text from a text frame while filtering out strikethrough text.
        
        Args:
            text_frame: A text frame object from python-pptx
            
        Returns:
            Extracted text with strikethrough content removed
        """
        text_parts = []
        try:
            for paragraph in text_frame.paragraphs:
                para_text = []
                for run in paragraph.runs:
                    # Skip strikethrough text using XML-based detection
                    if self._has_strikethrough(run):
                        continue
                    if run.text:
                        para_text.append(run.text)
                if para_text:
                    text_parts.append("".join(para_text))
        except Exception:
            # Fallback: try to get raw text if runs iteration fails
            if hasattr(text_frame, "text"):
                return text_frame.text.strip()
        return "\n".join(text_parts).strip()
    
    def _get_theme_hex(self, theme_color):
         
        theme_colors = {
            'ACCENT_1': '4472C4',
            'ACCENT_2': 'ED7D31',
            'ACCENT_3': 'A5A5A5',
            'ACCENT_4': 'FFC000',
            'ACCENT_5': '5B9BD5',
            'ACCENT_6': '70AD47',
            'HYPERLINK': '0563C1',
            'FOLLOWED_HYPERLINK': '954F72',
            'DARK_1': '000000',
            'LIGHT_1': 'FFFFFF',
            'DARK_2': '44546A',
            'LIGHT_2': 'E7E6E6',
            'BACKGROUND_1': 'FFFFFF',
            'TEXT_1': '000000',
            'BACKGROUND_2': 'E7E6E6',
            'TEXT_2': '44546A'
        }
        return theme_colors.get(str(theme_color.name), None)
    def _get_shape_fill_color(self, shape_fill):
        result = None
        try:
            if shape_fill.type == MSO_FILL_TYPE.SOLID:
                fore_color = shape_fill.fore_color

                if fore_color.type == MSO_COLOR_TYPE.RGB:
                    result = fore_color.rgb

                elif fore_color.type == MSO_COLOR_TYPE.SCHEME:
                    result = self._get_theme_hex(fore_color.theme_color)
        except Exception as e:
            print(f"Warning: Failed to get shape fill color: {e}")
        return result
    def _white_on_white(self, run, shape_color):
        """
        Check if a run has white text on a white shape background.
        
        Args:
            run: A run object from python-pptx paragraph
            shape_color: Hex color string of the shape background (e.g., 'FFFFFF')
        Returns:
            True if the run has white text and the shape background is also white
        """
        try:
            font_color = None
            if run.font and run.font.color:
                 
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    font_color = run.font.color.rgb
                elif run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                    font_color = self._get_theme_hex(run.font.color.theme_color)
            if font_color is not None and shape_color is not None:
                return (font_color == shape_color)
        except Exception as e:
            print(f"[DEBUG] Failed to determine white-on-white: {e}")
        
        return False
    def _get_slidemaster_theme(self, slide_master):
        sm_elm = slide_master._element
        ns = sm_elm.nsmap

        bg_ref = sm_elm.xpath(".//p:bgRef")
        if bg_ref:
            idx = bg_ref[0].get("idx")
        else:
            return None
        scheme_clr = bg_ref[0].xpath(".//a:schemeClr", namespaces = ns)
        if scheme_clr:
            scheme_val = scheme_clr[0].get("val")
        else:
            return None
        clr_map_node = sm_elm.xpath(".//p:clrMap")
        if not clr_map_node:
            return None
        mapped_val =   clr_map_node[0].get(scheme_val)
        if mapped_val:
            return mapped_val
        
        return None
    
    def _get_theme_hex_from_slidemaster(self, theme_color: str):
        if theme_color is None:
            return ""
        theme_colors = {
            'accent1': '4472C4',
            'accent2': 'ED7D31',
            'accent3': 'A5A5A5',
            'accent4': 'FFC000',
            'accent5': '5B9BD5',
            'accent6': '70AD47',
            'hlink': '0563C1',
            'folHlink': '954F72',
            'dk1': '000000',
            'lt1': 'FFFFFF',
            'dk2': '44546A',
            'lt2': 'E7E6E6'
        }
        theme_value = theme_colors[theme_color]
        return theme_value
          
        
    def _extract_text_with_line_bboxes(
        self, 
        shape: BaseShape, slide_master: None
    ) -> Tuple[str, List[LineBBox]]:
        """
        Extract text from a shape along with bounding boxes for each paragraph.
        
        This method calculates accurate bounding boxes for each paragraph
        based on font size, line spacing, space_after, margins, and text wrapping.
        
        Args:
            shape: A shape containing text (must have text_frame attribute)
            
        Returns:
            Tuple of (extracted_text, list of LineBBox objects)
        """
        if not hasattr(shape, "text_frame"):
            text = shape.text.strip() if hasattr(shape, "text") else ""
            if text:
                bbox = BoundingBox.from_shape(shape)
                return text, [LineBBox(text=text, bbox=bbox)]
            return "", []
        
        text_frame = shape.text_frame
        shape_bbox = BoundingBox.from_shape(shape)
        
        # Get margins
        margin_top = text_frame.margin_top or 0
        margin_bottom = text_frame.margin_bottom or 0
        margin_left = text_frame.margin_left or 0
        margin_right = text_frame.margin_right or 0
        
        # Available width for text (used for wrap calculation)
        text_width = shape_bbox.width - margin_left - margin_right
        
        shape_fill = shape.fill if hasattr(shape, "fill") else None
        shape_color = None
        if shape_fill is not None:
            shape_color = self._get_shape_fill_color(shape_fill)
        if shape_color is None:
            shape_value = self._get_slidemaster_theme(slide_master)
            shape_color = self._get_theme_hex_from_slidemaster(shape_value)
        # Extract paragraphs with their text (filtering strikethrough)
        paragraphs_data = []  # List of (text, font_size, line_spacing, space_after, num_wrapped_lines)
        try:
            for paragraph in text_frame.paragraphs:
                para_text = []
                for run in paragraph.runs:
                     
                    

                    if self._has_strikethrough(run):
                        continue
                    if self._white_on_white(run, shape_color):
                        continue
                    if run.text:
                        para_text.append(run.text)
                
                if not para_text:
                    continue
                    
                text = "".join(para_text)
                
                # Get font size
                font_size = None
                for run in paragraph.runs:
                    if run.font and run.font.size:
                        font_size = run.font.size
                        break
                if font_size is None:
                    font_size = Pt(10)  # Default 10pt
                
                # Get line spacing (default 1.0)
                line_spacing = paragraph.line_spacing if paragraph.line_spacing else 1.0
                
                # Get space after
                space_after = paragraph.space_after if paragraph.space_after else 0
                
                # Estimate number of wrapped lines
                # Character width is approximately 80% of font size for CJK characters
                char_width = font_size * 0.8
                chars_per_line = int(text_width / char_width) if char_width > 0 else 50
                num_wrapped_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
                
                paragraphs_data.append({
                    'text': text,
                    'font_size': font_size,
                    'line_spacing': line_spacing,
                    'space_after': space_after,
                    'num_wrapped_lines': num_wrapped_lines
                })
                
        except Exception as e:
            print(f"[DEBUG] Error in extract text with line bboxes : {e}")
            if hasattr(text_frame, "text") and text_frame.text.strip():
                return text_frame.text.strip(), [LineBBox(text=text_frame.text.strip(), bbox=shape_bbox)]
            
            return "", []
        
        if not paragraphs_data:
            return "", []
        
        # Calculate total content height
        total_content_height = 0
        for para in paragraphs_data:
            line_height = para['font_size'] * para['line_spacing']
            para_height = line_height * para['num_wrapped_lines'] + para['space_after']
            total_content_height += para_height
        
        # Available height for content
        available_height = shape_bbox.height - margin_top - margin_bottom
        
        # Scale factor if content exceeds available space (or if we need to stretch)
        scale_factor = available_height / total_content_height if total_content_height > 0 else 1.0
        
        # Build line bboxes
        line_bboxes = []
        current_y = shape_bbox.y + margin_top
        
        for para in paragraphs_data:
            line_height = para['font_size'] * para['line_spacing']
            para_height = line_height * para['num_wrapped_lines']
            space_after = para['space_after']
            
            # Apply scale factor
            scaled_para_height = para_height * scale_factor
            scaled_space_after = space_after * scale_factor
            
            line_bbox = BoundingBox(
                x=shape_bbox.x + margin_left,
                y=int(current_y),
                width=int(text_width),
                height=int(scaled_para_height)
            )
            line_bboxes.append(LineBBox(text=para['text'], bbox=line_bbox))
            
            current_y += scaled_para_height + scaled_space_after
        
        full_text = "\n".join(para['text'] for para in paragraphs_data)
        return full_text, line_bboxes

    def _extract_text_from_shape(self, shape: BaseShape) -> str:
        """Extract text content from a shape, filtering out strikethrough text."""
        if hasattr(shape, "text_frame"):
            return self._extract_text_from_text_frame(shape.text_frame)
        elif hasattr(shape, "text"):
            return shape.text.strip()
        return ""

    def _get_cell_bbox(self, table_shape: GraphicFrame, row_idx: int, col_idx: int) -> Optional[BoundingBox]:
        """
        Calculate the bounding box of a specific table cell.
        
        Args:
            table_shape: The GraphicFrame containing the table
            row_idx: Row index (0-based)
            col_idx: Column index (0-based)
            
        Returns:
            BoundingBox of the cell or None if cannot be calculated
        """
        try:
            table = table_shape.table
            # Get table position
            table_left = table_shape.left or 0
            table_top = table_shape.top or 0
            
            # Calculate cumulative row heights
            row_top = table_top
            for i in range(row_idx):
                row_top += table.rows[i].height or 0
            
            row_height = table.rows[row_idx].height or 0
            
            # Calculate cumulative column widths
            col_left = table_left
            for i in range(col_idx):
                col_left += table.columns[i].width or 0
            
            col_width = table.columns[col_idx].width or 0
            
            return BoundingBox(
                x=col_left,
                y=row_top,
                width=col_width,
                height=row_height
            )
        except Exception:
            return None

    def _is_cell_deprecated(
        self, 
        cell_bbox: BoundingBox, 
        cell_z_index: int,
        overlaying_elements: List[SlideElement],
        overlap_threshold: float = 0.2
    ) -> bool:
        """
        Check if a table cell is deprecated (covered by overlaying shapes).
        
        Args:
            cell_bbox: Bounding box of the cell
            cell_z_index: Z-index of the table containing this cell
            overlaying_elements: List of elements that can overlay content
            overlap_threshold: Minimum overlap ratio to consider cell deprecated
            
        Returns:
            True if cell is covered by overlaying shapes
        """
        if not overlaying_elements:
            return False

        # area-based overlap relative to the cell
        try:
            cell_area = max(1, cell_bbox.width * cell_bbox.height)
        except Exception:
            cell_area = 1

        # Center point of the cell for point-in-overlay heuristics
        center_x = int(cell_bbox.x + cell_bbox.width / 2)
        center_y = int(cell_bbox.y + cell_bbox.height / 2)

        for overlay in overlaying_elements:
            try:
                if not overlay or not overlay.bbox:
                    continue

                ov_bbox = overlay.bbox
                overlap_ratio = cell_bbox.get_overlap_ratio(ov_bbox)

                # If an overlay substantially covers the cell, it's deprecated
                if overlap_ratio >= overlap_threshold:
                    # Prefer obvious visual overlays regardless of z-index
                    if overlay.element_type in (ElementType.IMAGE, ElementType.GROUP):
                        return True

                    # Connectors should only deprecate if they pass through the cell center
                    if overlay.element_type == ElementType.CONNECTOR:
                        if ov_bbox.contains_point(center_x, center_y) or overlap_ratio > 0.5:
                            return True
                        continue

                    # If overlay is a shape/textbox/placeholder and appears as an overlaying shape
                    # (e.g. filled shape, opaque textbox), use existing overlay detection
                    if overlay.element_type in (ElementType.SHAPE, ElementType.TEXTBOX, ElementType.PLACEHOLDER):
                        if self._is_overlaying_shape(overlay):
                            return True

                        # If the overlay carries visible text or covers a large portion, consider deprecated
                        if overlay.content and overlay.content.strip():
                            return True
                        if overlap_ratio > 0.3:
                            return True

                    # Conservative fallback based on z-index when nothing else indicates transparency
                    if overlay.z_index >= cell_z_index and overlap_ratio >= overlap_threshold:
                        return True

                # Even if overlap_ratio is small, if the overlay covers the geometric center
                # of the cell and is visually overlaying, treat as deprecated
                if ov_bbox.contains_point(center_x, center_y):
                    if overlay.z_index >= cell_z_index or self._is_overlaying_shape(overlay):
                        return True

            except Exception:
                # ignore overlays that cause unexpected errors
                continue

        return False

    def _extract_cell_text(
        self,
        cell,
        table_shape: GraphicFrame,
        row_idx: int,
        col_idx: int,
        table_z_index: int,
        overlaying_elements: List[SlideElement]
    ) -> str:
        """
        Extract text from a table cell, checking for deprecation.
        
        Args:
            cell: The cell object to extract text from
            table_shape: The GraphicFrame containing the table
            row_idx: Row index of the cell
            col_idx: Column index of the cell
            table_z_index: Z-index of the table
            overlaying_elements: List of elements that can overlay content
            
        Returns:
            Extracted cell text or empty string if deprecated
        """
        is_deprecated = False
        if table_shape is not None and overlaying_elements:
            cell_bbox = self._get_cell_bbox(table_shape, row_idx, col_idx)
            if cell_bbox:
                is_deprecated = self._is_cell_deprecated(
                    cell_bbox, table_z_index, overlaying_elements
                )
        
        if is_deprecated:
            return ""
        
        if hasattr(cell, "text_frame"):
            cell_text = self._extract_text_from_text_frame(cell.text_frame)
        else:
            cell_text = cell.text.strip() if hasattr(cell, "text") else ""
        
        cell_text = re.sub(r'\s+', ' ', cell_text)
        cell_text = cell_text.replace('|', '\\|')
        return cell_text

    def _build_flattened_table_grid(
        self, 
        table: Table,
        table_shape: GraphicFrame = None,
        table_z_index: int = 0,
        overlaying_elements: List[SlideElement] = None
    ) -> List[List[str]]:
        """
        Build a flattened 2D grid from a table with merged cells.
        
        For merged cells, the content from the merge origin is propagated
        to all cells within the merge span. This ensures that when a cell
        spans multiple rows/columns, all grid positions covered by that
        cell contain the same value.
        
        Args:
            table: The Table object to process
            table_shape: Optional GraphicFrame for calculating cell positions
            table_z_index: Z-index of the table for overlay comparison
            overlaying_elements: Optional list of elements that can overlay content
            
        Returns:
            2D list of strings representing the flattened table grid
        """
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        
        # Initialize grid with empty strings
        grid = [["" for _ in range(num_cols)] for _ in range(num_rows)]
        
        # Track which cells have been filled (to handle overlapping merges)
        filled = [[False for _ in range(num_cols)] for _ in range(num_rows)]
        
        for row_idx in range(num_rows):
            for col_idx in range(num_cols):
                # Skip if already filled by a previous merge
                if filled[row_idx][col_idx]:
                    continue
                
                try:
                    cell = table.cell(row_idx, col_idx)
                except Exception:
                    continue
                
                # Get cell content (with deprecation check)
                cell_text = self._extract_cell_text(
                    cell, table_shape, row_idx, col_idx, 
                    table_z_index, overlaying_elements
                )
                
                # Check if this is a merge origin
                try:
                    is_merge_origin = cell.is_merge_origin
                    span_h = cell.span_height if is_merge_origin else 1
                    span_w = cell.span_width if is_merge_origin else 1
                except Exception:
                    is_merge_origin = False
                    span_h = 1
                    span_w = 1
                
                if is_merge_origin and (span_h > 1 or span_w > 1):
                    # Fill all cells in the span with the origin's content
                    for r in range(row_idx, min(row_idx + span_h, num_rows)):
                        for c in range(col_idx, min(col_idx + span_w, num_cols)):
                            grid[r][c] = cell_text
                            filled[r][c] = True
                else:
                    grid[row_idx][col_idx] = cell_text
                    filled[row_idx][col_idx] = True
        
        return grid

    def _extract_table_content(
        self, 
        table: Table,
        table_shape: GraphicFrame = None,
        table_z_index: int = 0,
        overlaying_elements: List[SlideElement] = None
    ) -> Tuple[List[List[str]], str, Dict[str, Any]]:
        """
        Extract table content, format as markdown, and compute table metadata.
        
        Cells covered by overlaying shapes (higher z-index) are marked as deprecated
        and their content is replaced with empty string.
        
        Merged cells are flattened: the content from the merge origin cell is
        propagated to all cells within the merge span.

        Args:
            table: The Table object to extract
            table_shape: Optional GraphicFrame for calculating cell positions
            table_z_index: Z-index of the table for overlay comparison
            overlaying_elements: Optional list of elements that can overlay content

        Returns:
            Tuple of (table_data as 2D list, markdown string, table_metadata dict)
            
        table_metadata includes:
            - column_count: int
            - column_widths: List[float] (normalized 0-1)
            - first_row: List[str]
            - total_width: int (in EMUs)
        """
        column_widths_emu = []
        
        # Extract column widths from table columns
        try:
            total_width = sum(col.width for col in table.columns)
            column_widths_emu = [col.width for col in table.columns]
        except Exception:
            total_width = 0
            column_widths_emu = []
        
        # Build flattened grid that handles merged cells
        rows = self._build_flattened_table_grid(
            table, table_shape, table_z_index, overlaying_elements
        )

        if not rows:
            return [], "", {"column_count": 0, "column_widths": [], "first_row": [], "total_width": 0}
        
        max_cols = max(len(row) for row in rows) if rows else 0
        for row in rows:
            while len(row) < max_cols:
                row.append("")
        
        # Normalize column widths to 0-1 range
        if total_width > 0 and column_widths_emu:
            column_widths_normalized = [w / total_width for w in column_widths_emu]
            # Pad if fewer widths than columns
            while len(column_widths_normalized) < max_cols:
                column_widths_normalized.append(1.0 / max_cols)
        else:
            # Default to equal widths
            column_widths_normalized = [1.0 / max_cols] * max_cols if max_cols > 0 else []
        
        # Build markdown
        markdown_lines = []
        if rows:
            markdown_lines.append("| " + " | ".join(rows[0]) + " |")
            markdown_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
            for row in rows[1:]:
                markdown_lines.append("| " + " | ".join(row) + " |")

        table_metadata = {
            "column_count": max_cols,
            "column_widths": column_widths_normalized,
            "first_row": rows[0] if rows else [],
            "total_width": total_width,
        }
        
        return rows, "\n".join(markdown_lines), table_metadata

    def _extract_image_from_shape(self, shape: Picture) -> Tuple[Optional[bytes], str]:
        """
        Extract image data from a Picture shape.

        Returns:
            Tuple of (image_bytes, extension)
        """
        try:
            image = shape.image
            image_bytes = image.blob
            content_type = image.content_type
            ext_map = {
                "image/png": "png",
                "image/jpeg": "jpg",
                "image/gif": "gif",
                "image/bmp": "bmp",
                "image/tiff": "tiff",
                "image/x-wmf": "wmf",
                "image/x-emf": "emf",
            }
            extension = ext_map.get(content_type, "png")

            return image_bytes, extension
        except Exception as e:
            print(f"Warning: Failed to extract image: {e}")
            return None, "png"

    def _capture_group_as_image(
        self, slide_index: int, bbox: BoundingBox
    ) -> Tuple[Optional[bytes], str]:
        """
        Capture a region of a slide as an image (for diagrams/groups).

        Uses LibreOffice to export the slide, then crops the specified region.

        Args:
            slide_index: Index of the slide (0-based)
            bbox: Bounding box of the region to capture

        Returns:
            Tuple of (im4-1.1	　　　はじめに　INTRODUCTION
		規定対象　Scope of Document	
		関連仕様書　Relevant Specificationsage_bytes, extension) or (None, "png") if failed
        """
        try:
            slide_image_path = self._get_exported_slide_path(slide_index - 1)
            if not slide_image_path:
                print(f"[DEBUG] No image slide path")
                return None, "png"
            image_bytes = self._crop_region_from_slide(slide_image_path, bbox)
            if image_bytes:
                return image_bytes, "png"
            print("[DEBUG] Can not get cropped image")
            return None, "png"
        except Exception as e:
            print(f"Warning: Failed to capture group as image: {e}")
            return None, "png"

    def _get_shape_type_name(self, shape: BaseShape) -> str:
        """Get a human-readable name for the shape type."""
        try:
            if isinstance(shape, AutoShape):
                if hasattr(shape, "auto_shape_type") and shape.auto_shape_type:
                    return str(shape.auto_shape_type).split(".")[-1]
            if hasattr(shape, "shape_type"):
                return str(shape.shape_type).split(".")[-1]
        except:
            pass
        return "UNKNOWN"
    def _get_ignored_images(self):
        ignored_paths = get_image_urls(self._repository_id)
        ignored_images = []
        for image_path in ignored_paths:
            url = get_signed_url(image_path)
            try:
                response = requests.get(url)
                image_bytes  = response.content
        
                ignored_images.append(image_bytes)
            except Exception as e:
                print(f"[DEBUG] Failed in getting image byte content  : {e}")
                
        return ignored_images
    # def _get_local_ignored_images(self):
    #     ignore_dir = '/home/namvt27/Project_TMC/ai-backend/functional_test/main/app/reader/run_script/ignored_images'
        
    #     ignored_images = []

    #     for image in os.listdir(ignore_dir):
    #         image_path = os.path.join(ignore_dir, image)
    #         with open(image_path, "rb") as f:
    #             image_bytes = f.read()

    #         ignored_images.append(image_bytes)
            
        
    #     return ignored_images
    def _check_ignored_image(self, image_bytes: io.BytesIO) -> bool:
        image = Image.open(io.BytesIO(image_bytes))
        
        for ignored_image in self.ignored_images:
            target_hash = phash(Image.open(io.BytesIO(ignored_image)))
            current_hash = phash(image)
            sim = similarity(current_hash, target_hash)
            print(f"[DEBUG] Object similarity : {sim}")
            if sim >= 0.7:
                print(f"[DEBUG] image has been ignored.")
                return True
        return False
        
    def _erase_grouped_image_covered(self, merged_overlaying: list[SlideElement], cluster_elements: list[SlideElement]):
        remaining_elements = []
        to_remove = []
        for element in cluster_elements:
            covering_overlays = [
            ov for ov in merged_overlaying
            if ov.z_index > element.z_index and element.bbox.get_overlap_ratio(ov.bbox) > 0.6
        ]
            if not covering_overlays:
                remaining_elements.append(element)
            else:
                covering_overlays.append(element)
                to_remove.extend(covering_overlays)
        
        result = [p for p in remaining_elements if p not in to_remove]
      
        return result
                

    # =========================================================================
    # STEP 1: Build Slide Structure
    # =========================================================================
    def _extract_element(
        self, shape: BaseShape, slide_index: int, z_index: int = 0, parent_bbox: Optional[BoundingBox] = None, slide_master=None
    ) -> Optional[SlideElement]:
        """
        Extract a single element from a shape with z-index (stacking order).

        Returns None if the element should be skipped (empty invisible shape).
        """

        element_type = self._get_element_type(shape)
        bbox = BoundingBox.from_shape(shape)

        element = SlideElement(
            element_type=element_type,
            bbox=bbox,
            raw_shape=shape,
            shape_type_name=self._get_shape_type_name(shape),
            slide_index=slide_index,
            z_index=z_index,
        )
        if element_type == ElementType.GROUP:
            group_shape: GroupShape = shape
            print("[DEBUG] GROUp shape detected")
            for child_shape in group_shape.shapes:
                child_element =  self._extract_element(child_shape, slide_index, z_index, bbox, slide_master)
                if child_element is not None:
                    element.children.append(child_element)
            # if self.get_image_content:
            #     try:
            #         print(f"[DEBUG] Group image printed !")
            #         image_bytes, ext = self._capture_group_as_image(slide_index, bbox)
                    
            #         if image_bytes:
            #             path, hash_val = self._save_image(image_bytes, ext)
            #             element.image_path = path
            #             element.image_hash = hash_val
            #             element.content = f"{IMAGE_PATH_START_MARKER} {path}|{hash_val} {IMAGE_PATH_END_MARKER}"
            #         else:
                    
            #             text = self._extract_text_from_shape(shape)
            #             if text:
            #                 element.content = f"{IMAGE_CONTENT_START_MARKER} [DIAGRAM] {text} {IMAGE_CONTENT_END_MARKER}"
            #             else:
            #                 element.content = f"{IMAGE_CONTENT_START_MARKER} [DIAGRAM] {IMAGE_CONTENT_END_MARKER}"
            #     except Exception as e:
            #         print(f"[DEBUG] Failed to capture group image : {type(e).__name__} : {e}")
            # else:
            #     element.content = f"{IMAGE_PATH_START_MARKER} DIAGRAM {IMAGE_PATH_END_MARKER}"

        elif element_type == ElementType.IMAGE:
            if self.get_image_content:
                image_bytes, ext = self._extract_image_from_shape(shape)
                is_ignored = self._check_ignored_image(image_bytes)
                if is_ignored:
                    return None
                if image_bytes:
                    path, hash_val = self._save_image(image_bytes, ext)
                    element.image_path = path
                    element.image_hash = hash_val
                    element.content = f"{IMAGE_PATH_START_MARKER} {path}|{hash_val} {IMAGE_PATH_END_MARKER}"
                else:
                    element.content = f"{IMAGE_PATH_START_MARKER} MISSING_IMAGE {IMAGE_PATH_END_MARKER}"
            else:
                element.content = f"{IMAGE_PATH_START_MARKER} IMAGE {IMAGE_PATH_END_MARKER}"

        elif element_type == ElementType.TABLE:
            if self.get_table_content and isinstance(shape, GraphicFrame) and shape.has_table:
                table_data, markdown, table_metadata = self._extract_table_content(shape.table)
                element.table_data = table_data
                element.content = f"{TABLE_CONTENT_START_MARKER} {markdown} {TABLE_CONTENT_END_MARKER}"
                # Populate table metadata for continuation detection
                element.table_column_count = table_metadata["column_count"]
                element.table_column_widths = table_metadata["column_widths"]
                element.table_first_row = table_metadata["first_row"]
                # Assign unique table ID
                self._table_id_counter += 1
                element.table_id = f"table_{element.slide_index}_{self._table_id_counter}"
            else:
                element.content = f"{TABLE_CONTENT_START_MARKER} TABLE {TABLE_CONTENT_END_MARKER}"

        elif element_type in (ElementType.TEXTBOX, ElementType.PLACEHOLDER):
            # Extract text with line-level bounding boxes for accurate deprecation detection
            content, line_bboxes = self._extract_text_with_line_bboxes(shape, slide_master)
            element.content = content
            element.line_bboxes = line_bboxes if line_bboxes else None
            if not element.content.strip():
                return None

        elif element_type == ElementType.SHAPE:
            # pass
            print("[DEBUG] Shape detected ")
            if self.get_shape_content:
                text = self._extract_text_from_shape(shape)
                if text:
                    element.content = f"{SHAPE_CONTENT_START_MARKER} {text} {SHAPE_CONTENT_END_MARKER}"
                else:
                    element.content = "" # Keep element for overlay detection, but no text
            else:
                element.content = f"{SHAPE_CONTENT_START_MARKER} SHAPE {SHAPE_CONTENT_END_MARKER}"

        elif element_type == ElementType.CONNECTOR:
            element.content = "" 

        return element

    def _is_overlaying_shape(self, element: SlideElement) -> bool:
        """
        Check if an element is a shape/diagram that can overlay and deprecate text.

        Overlaying shapes include:
        - SHAPE elements (rectangles, ovals, etc.) with significant fill
        - GROUP elements (diagrams)
        - IMAGE elements
        
        NOT overlaying:
        - TEXTBOX/PLACEHOLDER (text content)
        - CONNECTOR (lines/arrows - too thin to hide content)
        - Shapes without fill (transparent)

        Args:
            element: Element to check

        Returns:
            True if this element can overlay and deprecate underlying content
        """
        # Groups and images can definitely overlay content
        if element.element_type in (ElementType.GROUP, ElementType.IMAGE, ElementType.CONNECTOR):
            return True
            

        # Shapes can overlay if they have fill
        if element.element_type == ElementType.SHAPE:
            try:
                shape = element.raw_shape
                if hasattr(shape, 'fill'):
                    fill_type = shape.fill.type
                    # Check if shape has solid or patterned fill (not transparent)
                    if fill_type in (MSO_FILL_TYPE.SOLID, MSO_FILL_TYPE.PATTERNED, 
                                     MSO_FILL_TYPE.PICTURE, MSO_FILL_TYPE.GRADIENT):
                        return True
                    # Background fill might also overlay
                    if fill_type == MSO_FILL_TYPE.BACKGROUND:
                        return True
            except Exception:
                # If we can't determine fill, assume it could overlay
                return True
        
        return False
    def _is_overlaying_cluster(self, element: SlideElement) -> bool:
        # Groups and images can definitely overlay content
        if element.element_type in (ElementType.GROUP, ElementType.IMAGE, ElementType.CONNECTOR, ElementType.TEXTBOX):
            return True
            

        # Shapes can overlay if they have fill
        if element.element_type == ElementType.SHAPE:
            try:
                shape = element.raw_shape
                if hasattr(shape, 'fill'):
                    fill_type = shape.fill.type
                    # Check if shape has solid or patterned fill (not transparent)
                    if fill_type in (MSO_FILL_TYPE.SOLID, MSO_FILL_TYPE.PATTERNED, 
                                     MSO_FILL_TYPE.PICTURE, MSO_FILL_TYPE.GRADIENT):
                        return True
                    # Background fill might also overlay
                    if fill_type == MSO_FILL_TYPE.BACKGROUND:
                        return True
            except Exception:
                # If we can't determine fill, assume it could overlay
                return True
        
        return False
    
    def _is_filterable_element(self, element: SlideElement) -> bool:
        """
        Check if an element is a content element that can be deprecated/filtered.

        Filterable elements:
        - TEXTBOX
        - PLACEHOLDER
        - TABLE (has text in cells)
        - IMAGE (can be covered)
        - SHAPE (can be covered)

        Args:
            element: Element to check

        Returns:
            True if element contains content that could be deprecated by overlays
        """
        return element.element_type in (
            ElementType.TEXTBOX, 
            ElementType.PLACEHOLDER, 
            ElementType.TABLE,
            ElementType.IMAGE,
            ElementType.SHAPE
        )
    def _check_erased_image(self, element: SlideElement, covering_overlays: List[SlideElement]):
        for obj in covering_overlays:
            if element.bbox.get_overlap_ratio(obj.bbox) > 0.9:
                return True
        return False
    def _filter_overlapped_content(
        self,
        element: SlideElement,
        overlaying_elements: List[SlideElement]
    ) -> Tuple[str, Optional[List[LineBBox]]]:
        """
        Filter out text content from regions covered by overlaying shapes.
        
        Uses actual line bounding boxes (if available) for accurate detection.
        Falls back to estimated positions if line_bboxes is not available.
        
        Strategy:
        1. Use actual line bounding boxes from element.line_bboxes
        2. For each line, check if its bbox is covered by any overlay
        3. Remove lines that are significantly covered
        
        Args:
            element: Text element to filter
            overlaying_elements: List of overlaying shapes
            
        Returns:   
            Tuple of (filtered_content, updated_line_bboxes)
        """
        if not element.content or not overlaying_elements:
            return element.content, element.line_bboxes
        
        # Get overlays that cover this element with higher z-index
        covering_overlays = [
            ov for ov in overlaying_elements
            if ov.z_index > element.z_index and element.bbox.get_overlap_ratio(ov.bbox) > 0 
        ]
        if len(covering_overlays) > 0:
            if element.element_type in [ElementType.IMAGE, ElementType.GROUP, ElementType.SHAPE]:
                for covering_obj in covering_overlays:
                    if covering_obj.element_type in [ElementType.SHAPE, ElementType.CONNECTOR]:
                        if self._check_erased_image(element, covering_overlays):
                            return "", None
             
        if not covering_overlays:
            return element.content, element.line_bboxes
        
        content = element.content
        
        # Skip special content markers (only tables are handled separately)
        # Images and Shapes SHOULD be filtered if covered
        if TABLE_CONTENT_START_MARKER in content:
            return content, element.line_bboxes
        
        # Use actual line bboxes if available
        if element.line_bboxes:
            return self._filter_lines_with_bboxes(element.line_bboxes, covering_overlays)
        
        # Fallback: estimate line positions
        return self._filter_lines_estimated(element, covering_overlays)

    def _filter_lines_with_bboxes(
        self,
        line_bboxes: List[LineBBox],
        covering_overlays: List[SlideElement],
        vertical_threshold: float = 0.4,
        horizontal_threshold: float = 0.5
    ) -> Tuple[str, List[LineBBox]]:
        """
        Filter lines using actual bounding boxes.
        
        Args:
            line_bboxes: List of LineBBox objects with actual positions
            covering_overlays: List of overlaying elements
            vertical_threshold: Minimum vertical overlap ratio to consider covered
            horizontal_threshold: Minimum horizontal overlap ratio to consider covered
            
        Returns:
            Tuple of (filtered_content, filtered_line_bboxes)
        """
        filtered_line_bboxes = []
        
        for line_bbox in line_bboxes:
            is_covered = False
            
            for ov in covering_overlays:
                if line_bbox.is_covered_by(
                    ov.bbox, 
                    vertical_threshold=vertical_threshold,
                    horizontal_threshold=horizontal_threshold
                ):
                    is_covered = True
                    # Debug: uncomment to see which lines are filtered
                    # print(f"[FILTERED LINE] '{line_bbox.text[:40]}...' covered by overlay (actual bbox)")
                    break
            
            if not is_covered:
                filtered_line_bboxes.append(line_bbox)
        
        filtered_content = '\n'.join(lb.text for lb in filtered_line_bboxes)
        return filtered_content, filtered_line_bboxes

    def _filter_lines_estimated(
        self,
        element: SlideElement,
        covering_overlays: List[SlideElement],
        vertical_threshold: float = 0.4,
        horizontal_threshold: float = 0.5
    ) -> Tuple[str, None]:
        """
        Filter lines using estimated positions (fallback when no line_bboxes).
        
        Args:
            element: Text element with content
            covering_overlays: List of overlaying elements
            vertical_threshold: Minimum vertical overlap ratio to consider covered
            horizontal_threshold: Minimum horizontal overlap ratio to consider covered
            
        Returns:
            Tuple of (filtered_content, None)
        """
        content = element.content
        lines = content.split('\n')
        
        if len(lines) <= 1:
            return content, None
        
        # Estimate line positions
        total_height = element.bbox.height
        line_height = total_height / len(lines) if lines else total_height
        
        filtered_lines = []
        for idx, line in enumerate(lines):
            # Estimate this line's vertical position
            line_top = element.bbox.y + (idx * line_height)
            line_bottom = line_top + line_height
            
            # Check if this line is covered by any overlay
            is_covered = False
            for ov in covering_overlays:
                # Check vertical overlap
                vertical_overlap = max(0, min(line_bottom, ov.bbox.y2) - max(line_top, ov.bbox.y))
                v_ratio = vertical_overlap / line_height if line_height > 0 else 0
                
                if v_ratio >= vertical_threshold:
                    # Check horizontal overlap
                    h_overlap = max(0, min(element.bbox.x2, ov.bbox.x2) - max(element.bbox.x, ov.bbox.x))
                    h_ratio = h_overlap / element.bbox.width if element.bbox.width > 0 else 0
                    
                    if h_ratio >= horizontal_threshold:
                        is_covered = True
                        # Debug: uncomment to see which lines are filtered
                        # print(f"[FILTERED LINE] '{line[:40]}...' covered by overlay (estimated)")
                        break
            
            if not is_covered:
                filtered_lines.append(line)
        
        return '\n'.join(filtered_lines), None

    def _build_slide_structure(
        self, slide, slide_index: int, slide_master
    ) -> List[SlideElement]:
        """
        Build structure of elements in a slide.

        Elements are extracted with z-index (stacking order) and deprecated content detection.
        Empty invisible shapes are filtered during extraction.
        
        Deprecated content detection uses multiple strategies:
        1. Z-order based: Text covered by higher z-index shapes/diagrams is deprecated
        2. Line-level filtering: Remove specific lines covered by overlays
        3. Table cell-level: Individual cells covered by overlays are excluded
        4. Legacy region-based: Gray rectangular shapes mark deprecated regions
        """
        
        # Collect all slide elements (paragraph, image, table, ...)
        elements = []
        
        group_counter = 0
        z_indexes = []
        for z_idx, shape in enumerate(slide.shapes):
            # print(f"[DEBUG] Element shape : {self._get_element_type(shape)}")
            z_indexes.append(z_idx)
            if self._get_element_type(shape) == ElementType.GROUP:
                group_counter += 1
            element = self._extract_element(shape, slide_index, z_idx, None, slide_master)
            if element is None or (not element.content and element.element_type != ElementType.CONNECTOR):
                continue
            # if not element.content:
            #     continue
            elements.append(element)
            
        print(f"[DEBUG] Z-indexes of elements: {z_indexes}")
        
        # Collect overlaying shapes for deprecation detection
        overlaying_elements = [
            elem for elem in elements 
            if self._is_overlaying_shape(elem)
        ]
                
        if overlaying_elements:
            for elem in elements:
            # Re-process tables with overlay information for cell-level deprecation
                if elem.element_type == ElementType.TABLE and elem.raw_shape:
                    shape = elem.raw_shape
                    if isinstance(shape, GraphicFrame) and shape.has_table:
                        table_data, markdown, table_metadata = self._extract_table_content(
                            shape.table,
                            table_shape=shape,
                            table_z_index=elem.z_index,
                            overlaying_elements=overlaying_elements
                        )
                        elem.table_data = table_data
                        elem.content = f"{TABLE_CONTENT_START_MARKER} {markdown} {TABLE_CONTENT_END_MARKER}"
                        elem.table_column_count = table_metadata["column_count"]
                        elem.table_column_widths = table_metadata["column_widths"]
                        elem.table_first_row = table_metadata["first_row"]

                # Filter overlapped content (text, images, shapes)
                # Uses actual line bounding boxes or estimation
                if self._is_filterable_element(elem) and elem.element_type != ElementType.TABLE:
                    if elem.element_type == ElementType.IMAGE:
                        print("[DEBUG] Extracted image overlaying elements")
                    filtered_content, filtered_line_bboxes = self._filter_overlapped_content(elem, overlaying_elements)
                    if filtered_content == "erased element":
                        elem.content = ""
                        elem.line_bboxes = None
                    if filtered_content != elem.content:
                        elem.content = filtered_content
                        elem.line_bboxes = filtered_line_bboxes
                    

        
        # Remove deprecated elements from result
        # Also remove overlaying shapes that are just visual markers (shapes with "Delete" text)
        result = []
        for elem in elements:
            # if elem.is_deprecated:
            #     continue
            # Optionally filter out "Delete" marker shapes
            if elem.element_type == ElementType.SHAPE and self._is_overlaying_shape(elem):
                # Check if this is just a deprecation marker shape
                content_text = elem.content.lower()
                if 'delete' in content_text or '削除' in content_text:
                    # Debug: uncomment to see which Delete shapes are removed
                    # print(f"[FILTERED] Removing Delete marker shape: {elem.content[:50]}")
                    continue
            result.append(elem)
        
        result.sort()
        self._extracted_element = result
        return result

    # =========================================================================
    # STEP 2: Merge Elements
    # =========================================================================

    def _is_merge_candidate(self, element: SlideElement) -> bool:
        """
        Check if an element is a candidate for merging.

        Candidates (diagram components):
        - SHAPE elements (geometry shapes with fill)
        - CONNECTOR elements (arrows, lines)

        NOT candidates (standalone content):
        - TEXTBOX elements (text content - must remain standalone for header detection)
        - TABLE elements (standalone data)
        - GROUP elements (already merged)
        - IMAGE elements (standalone pictures)
        - PLACEHOLDER elements
        """
        if element.is_merged:
            return False
        if element.element_type == ElementType.GROUP:
            return True
        if element.element_type == ElementType.TABLE:
            return False
        if element.element_type == ElementType.PLACEHOLDER:
            return False
        if element.element_type == ElementType.TEXTBOX:
            for elem in self._extracted_element:
                if elem.bbox.get_overlap_ratio(element.bbox) == 1 and elem.element_type in [ElementType.IMAGE, ElementType.GROUP, ElementType.SHAPE]:
                    return True
            return False
        if element.element_type == ElementType.IMAGE:
            return True
        # if self._slide_width > 0:
        #     width_ratio = element.bbox.width / self._slide_width
        #     if width_ratio >= 0.5:
        #         return False

        return True

    def _should_merge(self, elem1: SlideElement, elem2: SlideElement) -> bool:
        """
        Determine if two elements should be merged based on proximity and alignment.
        """
        if not (self._is_merge_candidate(elem1) and self._is_merge_candidate(elem2)):
            return False
        vertical_gap = abs(elem2.bbox.y - elem1.bbox.y2)
        max_vertical_gap = self._slide_height * 0.1 if self._slide_height > 0 else Emu(Pt(50))
        horizontal_overlap = (
            max(0, min(elem1.bbox.x2, elem2.bbox.x2) - max(elem1.bbox.x, elem2.bbox.x))
        )
        if vertical_gap < max_vertical_gap:
            if horizontal_overlap > 0:
                return True
            horizontal_gap = abs(elem2.bbox.x - elem1.bbox.x2)
            max_horizontal_gap = self._slide_width * 0.05 if self._slide_width > 0 else Emu(Pt(30))
            if horizontal_gap < max_horizontal_gap:
                return True

        return False

    def _are_spatially_close(
        self,
        e1: SlideElement,
        e2: SlideElement,
        threshold_ratio: float = 0.08
    ) -> bool:
        """
        Check if two elements are spatially close based on centroid distance.

        Args:
            e1, e2: Elements to compare
            threshold_ratio: Maximum distance ratio (relative to slide dimensions)

        Returns:
            True if elements are close enough to be in same group
        """
        if self._slide_width <= 0 or self._slide_height <= 0:
            return False

        cx_dist = abs(e1.bbox.center_x - e2.bbox.center_x) / self._slide_width
        cy_dist = abs(e1.bbox.center_y - e2.bbox.center_y) / self._slide_height
        distance = (cx_dist * 2 + cy_dist * 2) ** 0.5
        return distance < threshold_ratio

    def _connector_touches_shape(
        self,
        connector: SlideElement,
        shape: SlideElement,
        tolerance_ratio: float = 0.02
    ) -> bool:
        """
        Check if a connector (arrow/line) touches or is near a shape.

        Args:
            connector: Connector element
            shape: Shape element
            tolerance_ratio: How close connector endpoints need to be

        Returns:
            True if connector appears to connect to the shape
        """
        if self._slide_width <= 0 or self._slide_height <= 0:
            return False
        tolerance_x = self._slide_width * tolerance_ratio
        tolerance_y = self._slide_height * tolerance_ratio
        conn_points = [
            (connector.bbox.x, connector.bbox.y),
            (connector.bbox.x2, connector.bbox.y2),
        ]
        for cx, cy in conn_points:
            near_x = (shape.bbox.x - tolerance_x <= cx <= shape.bbox.x2 + tolerance_x)
            near_y = (shape.bbox.y - tolerance_y <= cy <= shape.bbox.y2 + tolerance_y)
            if near_x and near_y:
                return True
        return False

    def _bfs_component(
        self,
        graph: Dict[int, set],
        start: int,
        visited: set
    ) -> List[int]:
        """
        BFS to find connected component in graph.

        Args:
            graph: Adjacency list representation
            start: Starting node
            visited: Set of already visited nodes (modified in place)

        Returns:
            List of node indices in the connected component
        """
        component = []
        queue = [start]
        visited.add(start)

        while queue:
            node = queue.pop(0)
            component.append(node)
            for neighbor in graph[node]:
                if neighbor not in visited:
                    visited.add(neighbor)
                    queue.append(neighbor)

        return component

    def _is_valid_merged_group(self, elements: List[SlideElement]) -> bool:
        """
        Check if a group of elements is valid for merging.

        A valid merged group must have at least one element that is NOT a textbox.
        Valid non-textbox elements:
        - SHAPE (geometry shapes like rectangles with fill, circles, etc.)
        - CONNECTOR (arrows, lines)
        - IMAGE

        If group contains ONLY textboxes, it should be split back to standalone elements.

        Args:
            elements: List of elements to check

        Returns:
            True if group is valid for merging (has at least one non-textbox element)
        """
        
        for elem in elements:
            if elem.element_type in (ElementType.SHAPE, ElementType.CONNECTOR, ElementType.IMAGE, ElementType.GROUP):
                return True
        return False

    def _merge_elements(self, elements: List[SlideElement]) -> List[SlideElement]:
        """
        Merge standalone first-level elements using hybrid approach.

        Strategy:
        1. DBSCAN spatial clustering for general diagrams
        2. Connector-aware graph merging for flowcharts with arrows
        3. Post-process to split overly wide groups (likely false positives)

        Merged groups are treated as diagrams.
        """
        if not elements:
            return elements

        result = self._merge_elements_dbscan(elements)
        # result = self._enhance_with_connectors(result)
        # result = self._split_wide_groups(result, max_width_ratio=0.65)

        return result
    def _check_table_contain(self, non_candidates: List[SlideElement], target_element: SlideElement) -> bool: 
        for element in non_candidates:
            if element.element_type == ElementType.TABLE:
                table_bbox = element.bbox
                table_x, table_y , table_x2, table_y2 = table_bbox.x, table_bbox.y, table_bbox.width, table_bbox.height
                target_bbox = target_element.bbox
                target_x, target_y, target_x2, target_y2 = target_bbox.x, target_bbox.y, target_bbox.width, target_bbox.height

                if target_x >= table_x and target_y >= table_y and target_x2 <= table_x2 and target_y2 <= table_y2:
                    formatted_content = ""
                    table_element_content = element.content
                    end_idx = table_element_content.find(TABLE_CONTENT_END_MARKER)
                    formatted_content = table_element_content[:end_idx]
                    formatted_content = formatted_content + "\n"+ target_element.content + table_element_content[end_idx:]
                    element.content = formatted_content
                    return True
        return False
    def _cleaned_covered_content(self, non_candidates: List[SlideElement], largest_element: SlideElement):
        for element in non_candidates:
            if element.bbox.get_overlap_ratio(largest_element.bbox) > 0.9:
                element.content = ""
                element.line_bboxes = None
    
    def _merge_elements_dbscan(self, elements: List[SlideElement]) -> List[SlideElement]:
        """
        Merge elements using DBSCAN spatial clustering.

        Args:
            elements: List of slide elements

        Returns:
            List with spatially clustered elements merged into groups
        """
        candidates = [e for e in elements if self._is_merge_candidate(e)]
        non_candidates = [e for e in elements if not self._is_merge_candidate(e)]
        if len(candidates) < 2:
            return elements
        centroids = np.array([[e.bbox.center_x, e.bbox.center_y] for e in candidates], dtype=np.float64)
        if self._slide_width > 0 and self._slide_height > 0:
            import math
            diagonal = math.sqrt(self._slide_width**2 + self._slide_height**2)
            centroids[:, 0] /= diagonal
            centroids[:, 1] /= diagonal
            
        clustering = DBSCAN(eps= 30, min_samples=1, metric='euclidean').fit(centroids)
        
        
        clusters = {}   
        for idx, label in enumerate(clustering.labels_):
            if label == -1:
                non_candidates.append(candidates[idx])
            else:
                if label not in clusters:
                    clusters[label] = []
                clusters[label].append(candidates[idx])
        print(f"[DEBUG] Identified clusters : {clustering.labels_}")
        
        merged_groups = []
        
        for cluster_elements in clusters.values():
            if len(cluster_elements) > 1:
                if self._is_valid_merged_group(cluster_elements) and self._is_valid_minimum(cluster_elements):
                    merged_overlying_elements = [p for p in cluster_elements if self._is_overlaying_cluster(p)]
                    print(f"Cluster elements before erasing: {len(cluster_elements)}")
                    cluster_elements= self._erase_grouped_image_covered(merged_overlying_elements, cluster_elements)
                    if len(cluster_elements) <= 1:
                        non_candidates.extend(cluster_elements)
                        continue
                    merged_element = self._create_merged_group(cluster_elements)
                    largest_element = max(cluster_elements, key=lambda e: e.bbox.width * e.bbox.height)
                    self._cleaned_covered_content(non_candidates, largest_element)
                    
                    
                    if not self._check_table_contain(non_candidates, merged_element):
                        merged_groups.append(merged_element)
                    # else:
                    #     merged_groups.append(self._create_merged_group(cluster_elements))
                else:
                    non_candidates.extend(cluster_elements)
            else:
                non_candidates.extend(cluster_elements)
        result = non_candidates + merged_groups
        result.sort(key=lambda e: (e.bbox.y, e.bbox.x))

        return result
    
    def _is_valid_minimum(self, cluster_elements: List[SlideElement]) -> bool:
        for element in cluster_elements:
            if element.element_type == ElementType.GROUP or len(cluster_elements) >= 4:
                return True
        return False
    def _split_wide_groups(
        self,
        elements: List[SlideElement],
        max_width_ratio: float = 0.65
    ) -> List[SlideElement]:
        """ 
        Post-process to split merged groups that are too wide.

        Overly wide groups are likely false positives (spanning unrelated content).

        Args:
            elements: Elements after initial merging
            max_width_ratio: Maximum allowed width as fraction of slide width

        Returns:
            Elements with overly wide groups split back into components
        """
        if self._slide_width <= 0:
            return elements

        result = []

        for elem in elements:
            if elem.is_merged and elem.bbox.width / self._slide_width > max_width_ratio:
                if elem.children:
                    result.extend(elem.children)
                else:
                    result.append(elem)
            else:
                result.append(elem)
        result.sort(key=lambda e: (e.bbox.y, e.bbox.x))

        return result

    def _create_merged_group(self, elements: List[SlideElement]) -> SlideElement:
        """Create a merged group from multiple elements, treating it as a diagram."""
        combined_bbox = elements[0].bbox
        for elem in elements[1:]:
            combined_bbox = combined_bbox.union(elem.bbox)
        slide_index = elements[0].slide_index
        print(f"Actually merged group slide_index : {slide_index}")
        merged = SlideElement(
            element_type=ElementType.GROUP,
            bbox=combined_bbox,
            children=elements,
            is_merged=True,
            slide_index=slide_index,
        )



        if self.get_image_content:
            
            print(f"[DEBUG] Created image for merged group.")
            image_bytes, ext = self._capture_group_as_image(slide_index, combined_bbox)
            if image_bytes:
                
                path, hash_val = self._save_image(image_bytes, ext)
                merged.image_path = path
                merged.image_hash = hash_val
                merged.content = f"{IMAGE_PATH_START_MARKER} {path}|{hash_val} {IMAGE_PATH_END_MARKER}"
            else:
                content_parts = []
                for elem in elements:
                    if elem.content:
                        content_parts.append(elem.content)
                inner_content = " | ".join(content_parts) if content_parts else "Merged Diagram"
                merged.content = f"{IMAGE_CONTENT_START_MARKER} [DIAGRAM] {inner_content} {IMAGE_CONTENT_END_MARKER}"
        else:
            merged.content = f"{IMAGE_PATH_START_MARKER} DIAGRAM {IMAGE_PATH_END_MARKER}"

        return merged

    # =========================================================================
    # STEP 4: Multi-Slide Table Merging
    # =========================================================================

    def _calculate_table_similarity(
        self,
        table1: SlideElement,
        table2: SlideElement,
    ) -> float:
        """
        Calculate similarity score (0-1) between two tables for continuation detection.

        Factors considered:
        - Column count match (required for high score)
        - Column width similarity
        - Position similarity (X position, width)
        - Header row match/absence

        Args:
            table1: First table (potential parent)
            table2: Second table (potential continuation)

        Returns:
            Similarity score from 0.0 to 1.0
        """
        if table1.element_type != ElementType.TABLE or table2.element_type != ElementType.TABLE:
            return 0.0
        
        score = 0.0
        
        # 1. Column count match (required - weight 0.4)
        col1 = table1.table_column_count or 0
        col2 = table2.table_column_count or 0
        
        if col1 == 0 or col2 == 0:
            return 0.0
        
        if col1 == col2:
            score += 0.4
        else:
            # Allow off-by-one for edge cases but penalize
            if abs(col1 - col2) == 1:
                score += 0.1
            else:
                return 0.0  # Column count mismatch is disqualifying
        
        # 2. Column width similarity (weight 0.25)
        widths1 = table1.table_column_widths or []
        widths2 = table2.table_column_widths or []
        
        if widths1 and widths2 and len(widths1) == len(widths2):
            width_diffs = [abs(w1 - w2) for w1, w2 in zip(widths1, widths2)]
            avg_diff = sum(width_diffs) / len(width_diffs) if width_diffs else 0
            # Score based on how close widths are (tolerance-based)
            if avg_diff <= self.table_column_width_tolerance:
                score += 0.25 * (1.0 - avg_diff / self.table_column_width_tolerance)
        
        # 3. Position similarity - X position and width (weight 0.2)
        if self._slide_width > 0:
            x_diff_ratio = abs(table1.bbox.x - table2.bbox.x) / self._slide_width
            width_diff_ratio = abs(table1.bbox.width - table2.bbox.width) / self._slide_width
            
            if x_diff_ratio < 0.05 and width_diff_ratio < 0.1:
                score += 0.2
            elif x_diff_ratio < 0.1 and width_diff_ratio < 0.15:
                score += 0.1
        
        # 4. Header row analysis (weight 0.15)
        first_row1 = table1.table_first_row or []
        first_row2 = table2.table_first_row or []
        
        if first_row1 and first_row2:
            # Check if headers are identical (continuation with repeated header)
            if first_row1 == first_row2:
                score += 0.15
            else:
                # Check if second table's first row looks like data (not a header)
                # Heuristic: if first row of table2 doesn't match table1's header pattern,
                # it's likely a data row continuation
                header_similarity = sum(1 for a, b in zip(first_row1, first_row2) if a == b) / max(len(first_row1), 1)
                if header_similarity < 0.3:
                    # Low similarity suggests table2 starts with data, not header
                    score += 0.15
                elif header_similarity > 0.7:
                    # High similarity suggests repeated header
                    score += 0.1
        
        return min(score, 1.0)

    def _is_table_continuation(
        self,
        parent_table: SlideElement,
        candidate_table: SlideElement,
    ) -> bool:
        """
        Determine if candidate_table is a continuation of parent_table.

        Requirements:
        - Tables must be on consecutive slides
        - Similarity score must exceed threshold
        - Column count must match (enforced in similarity calculation)

        Args:
            parent_table: The table that might be continued
            candidate_table: The table that might be a continuation

        Returns:
            True if candidate is a continuation of parent
        """
        # Must be on consecutive slides
        if candidate_table.slide_index != parent_table.slide_index + 1:
            return False
        
        # Already marked as continuation of another table
        if candidate_table.is_table_continuation:
            return False
        
        # Calculate similarity
        similarity = self._calculate_table_similarity(parent_table, candidate_table)
        
        return similarity >= self.table_similarity_threshold

    def _has_repeated_header(
        self,
        parent_table: SlideElement,
        continuation_table: SlideElement,
    ) -> bool:
        """
        Check if the continuation table has a repeated header row.

        Args:
            parent_table: The parent table
            continuation_table: The continuation table

        Returns:
            True if the first row of continuation matches parent's header
        """
        first_row1 = parent_table.table_first_row or []
        first_row2 = continuation_table.table_first_row or []
        
        if not first_row1 or not first_row2:
            return False
        
        if len(first_row1) != len(first_row2):
            return False
        
        # Check if headers match (allowing for minor variations)
        matches = sum(1 for a, b in zip(first_row1, first_row2) 
                      if a.strip().lower() == b.strip().lower())
        
        return matches >= len(first_row1) * 0.8  # 80% match threshold

    def _merge_table_data(
        self,
        parent_table: SlideElement,
        continuation_table: SlideElement,
    ) -> SlideElement:
        """
        Merge continuation table data into parent table.

        The continuation table's rows are appended to the parent table.
        If continuation has a repeated header, it is skipped.

        Args:
            parent_table: The parent table to merge into
            continuation_table: The continuation table to merge from

        Returns:
            Updated parent table with merged data
        """
        parent_data = parent_table.table_data or []
        continuation_data = continuation_table.table_data or []
        
        if not continuation_data:
            return parent_table
        
        # Check if we should skip the first row (repeated header)
        skip_first_row = self._has_repeated_header(parent_table, continuation_table)
        
        # Merge rows
        if skip_first_row and len(continuation_data) > 1:
            rows_to_add = continuation_data[1:]
        else:
            rows_to_add = continuation_data
        
        # Combine data
        merged_data = parent_data + rows_to_add
        
        # Rebuild markdown content
        if merged_data:
            max_cols = max(len(row) for row in merged_data)
            for row in merged_data:
                while len(row) < max_cols:
                    row.append("")
            
            markdown_lines = []
            markdown_lines.append("| " + " | ".join(merged_data[0]) + " |")
            markdown_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
            for row in merged_data[1:]:
                markdown_lines.append("| " + " | ".join(row) + " |")
            
            merged_markdown = "\n".join(markdown_lines)
        else:
            merged_markdown = ""
        
        # Track which slides this table spans
        merged_slides = parent_table.merged_from_slides or [parent_table.slide_index]
        if continuation_table.slide_index not in merged_slides:
            merged_slides.append(continuation_table.slide_index)
        
        # Create slide span comment
        if len(merged_slides) > 1:
            slide_span_comment = f"<!-- Table spans slides {merged_slides[0]+1}-{merged_slides[-1]+1} -->\n"
        else:
            slide_span_comment = ""
        
        print(f"[TABLE MERGING] {slide_span_comment.strip()}")
        
        # Update parent table
        parent_table.table_data = merged_data
        parent_table.content = f"{TABLE_CONTENT_START_MARKER} {merged_markdown} {TABLE_CONTENT_END_MARKER}"
        parent_table.merged_from_slides = merged_slides
        
        return parent_table

    def _merge_continuation_tables(
        self,
        elements: List[SlideElement],
    ) -> List[SlideElement]:
        """
        Post-process all elements to merge tables that span multiple slides.

        This method identifies tables that are continuations of tables from
        previous slides and merges them while preserving the element order
        for correct header tree structure.

        Algorithm:
        1. Collect all TABLE elements grouped by slide index
        2. For each slide (in order), check if tables are continuations of previous slide
        3. When a continuation is found:
           - Merge data into parent table
           - Mark continuation for removal
        4. Remove merged continuation tables from element list
        5. Preserve element order for correct tree building

        Args:
            elements: List of all slide elements (after header detection)

        Returns:
            List of elements with continuation tables merged into parent tables
        """
        if not self.merge_continuation_tables:
            return elements
        
        # Collect tables by slide index
        tables_by_slide: Dict[int, List[SlideElement]] = {}
        for elem in elements:
            if elem.element_type == ElementType.TABLE and elem.table_data:
                slide_idx = elem.slide_index
                if slide_idx not in tables_by_slide:
                    tables_by_slide[slide_idx] = []
                tables_by_slide[slide_idx].append(elem)
        
        if not tables_by_slide:
            return elements
        
        # Track which elements to remove (continuations that were merged)
        elements_to_remove: set = set()
        
        # Process slides in order
        sorted_slides = sorted(tables_by_slide.keys())
        
        for i, slide_idx in enumerate(sorted_slides):
            if i == 0:
                continue  # First slide can't have continuations
            
            prev_slide_idx = sorted_slides[i - 1]
            
            # Only consider consecutive slides
            if slide_idx != prev_slide_idx + 1:
                continue
            
            current_tables = tables_by_slide[slide_idx]
            prev_tables = tables_by_slide.get(prev_slide_idx, [])
            
            # For each table in current slide, check if it's a continuation
            for current_table in current_tables:
                if current_table.is_table_continuation:
                    continue  # Already processed
                
                # Find best matching parent table from previous slide
                best_parent = None
                best_similarity = 0.0
                
                for prev_table in prev_tables:
                    if self._is_table_continuation(prev_table, current_table):
                        similarity = self._calculate_table_similarity(prev_table, current_table)
                        if similarity > best_similarity:
                            best_similarity = similarity
                            best_parent = prev_table
                
                if best_parent:
                    # Merge current table into parent
                    self._merge_table_data(best_parent, current_table)
                    
                    # Mark current table as merged
                    current_table.is_table_continuation = True
                    current_table.continuation_of_table_id = best_parent.table_id
                    elements_to_remove.add(id(current_table))
        
        # Filter out merged continuation tables
        result = [elem for elem in elements if id(elem) not in elements_to_remove]
        
        return result

    def _group_elements_by_line(
        self, elements: List[SlideElement], y_tolerance_ratio: float = 0.01
    ) -> Dict[int, List[SlideElement]]:
        """
        Group elements that are on approximately the same horizontal line.

        Args:
            elements: List of slide elements
            y_tolerance_ratio: Y-coordinate tolerance as ratio of slide height

        Returns:
            Dictionary mapping line_index to list of elements on that line
        """
        if not elements or self._slide_height <= 0:
            return {}

        y_tolerance = self._slide_height * y_tolerance_ratio
        sorted_elements = sorted(elements, key=lambda e: e.bbox.y)

        lines = {}
        current_line_y = None
        line_index = 0

        for elem in sorted_elements:
            if current_line_y is None or abs(elem.bbox.y - current_line_y) > y_tolerance:
                line_index = len(lines)
                current_line_y = elem.bbox.y
                lines[line_index] = []

            lines[line_index].append(elem)

        return lines

    # =========================================================================
    # STEP 3: Header Detection and Splitting
    # =========================================================================

    async def _detect_and_split_headers_in_slide(
        self, elements: List[SlideElement], slide_idx: int
    ) -> List[SlideElement]:
        """
        Detect headers using GENERAL APPROACH: first line + pattern matching.

        Process:
        1. Split multi-line elements into individual line elements
        2. Group lines by horizontal position (y-coordinate)
        3. Combine content from elements on same line
        4. Apply GENERAL RULE:
           - First line (topmost) is ALWAYS a header candidate
           - Other lines become candidates IF they match header patterns
        5. Use LLM to validate all candidates
        6. Create separate header elements with is_header=True for validated ones

        This ensures slide titles/chapter headers at the top are always considered,
        regardless of their format.

        Args:
            elements: Elements from a single slide
            slide_idx: Slide index for LLM context

        Returns:
            Elements with headers split into separate textbox elements
        """
        split_elements = []

        for elem in elements:
            content = elem.content.strip()
            if not content:
                split_elements.append(elem)
                continue
            if content.startswith(IMAGE_PATH_START_MARKER) or content.startswith(TABLE_CONTENT_START_MARKER):
                split_elements.append(elem)
                continue
            text_content = content
            if SHAPE_CONTENT_START_MARKER in content:
                text_content = content.replace(SHAPE_CONTENT_START_MARKER, "").replace(
                    SHAPE_CONTENT_END_MARKER, ""
                ).strip()
            lines = text_content.split('\n')

            if len(lines) == 1:
                split_elements.append(elem)
            else:
                line_height = elem.bbox.height // len(lines) if len(lines) > 0 else elem.bbox.height

                for line_idx, line_text in enumerate(lines):
                    line_text = line_text.strip()
                    if not line_text:
                        continue
                    line_elem = SlideElement(
                        element_type=ElementType.TEXTBOX,
                        bbox=BoundingBox(
                            x=elem.bbox.x,
                            y=elem.bbox.y + (line_idx * line_height),
                            width=elem.bbox.width,
                            height=line_height
                        ),
                        content=line_text,
                        slide_index=elem.slide_index,
                        z_index=elem.z_index,
                    )
                    split_elements.append(line_elem)
       
        lines = self._group_elements_by_line(split_elements)
       

        
        candidates = []
        candidate_line_map = {}  

        # MAY NEED TO REMOVE THIS LOGIC BECAUSE IT ASSUMES FIRST LINE IS ALWAYS A HEADER
        first_line_idx = min(lines.keys()) if lines else None

        for line_idx, line_elements in sorted(lines.items()):
            
            line_elements.sort(key=lambda e: e.bbox.x)

            
            combined_content_parts = []
            for elem in line_elements:
                content = elem.content.strip()
                if not content:
                    continue
                
                if content.startswith(IMAGE_PATH_START_MARKER) or content.startswith(TABLE_CONTENT_START_MARKER):
                    continue

                if content:
                    combined_content_parts.append(content)

            if not combined_content_parts:
                continue

            
            combined_content = " ".join(combined_content_parts)
            
            # Determine if this line is in the top header area
            # This catches titles that might not be the very first element (e.g. obscured by logo/breadcrumb)
            line_y = line_elements[0].bbox.y if line_elements else 0
            is_top_area = False
            if self._slide_height > 0:
                is_top_area = line_y < (self._slide_height * 0.20)  # Top 20%

            is_first_line = (line_idx == first_line_idx)
            # We treat top area lines as "forced candidates" just like the first line
            # This allows titles to be detected even if they don't match standard patterns
            is_forced_candidate = is_first_line or is_top_area

            # --- Frequency-based marker filtering ---
            # If this forced candidate text appears too frequently across slides,
            # it's a repeated marker (company name, breadcrumb, logo label)   skip it.
            if is_forced_candidate and hasattr(self, '_marker_freq_map') and self._marker_scan_count > 0:
                norm_for_freq = normalize_line(combined_content).strip().lower()
                marker_count = self._marker_freq_map.get(norm_for_freq, 0)
                marker_ratio = marker_count / self._marker_scan_count
                if marker_ratio >= 0.40:
                    print(f"[MARKER_FILTER] Skipping repeated marker on slide {slide_idx}: "
                          f"'{combined_content[:60]}' (appears in {marker_count}/{self._marker_scan_count} = {marker_ratio:.0%} of slides)")
                    is_forced_candidate = False

            # Use heading_detector module for pattern matching
            # It handles bullet stripping, full-width normalization, and level detection
            header_match = None
            if not is_forced_candidate:
                # Check if it's a level 4+ pattern (treat as content, not header)
                level, pattern_name, match = detect_heading_level(combined_content)
                if level is None and pattern_name and match:
                    header_match = None
                elif level is not None and pattern_name and match:
                    # Detect heading level using the heading_detector module
                        header_match = {"level": level, "pattern": pattern_name, "match": match}
            
            if is_forced_candidate or header_match:
                candidate_idx = len(candidates)
                candidates.append(combined_content)
                candidate_line_map[candidate_idx] = {
                    'line_idx': line_idx,
                    'line_elements': line_elements,
                    'combined_content': combined_content,
                    'is_first_line': is_forced_candidate, # Pass True to indicate TOP position context
                    'y_position': line_y
                }
        

        validated_indices = set()
        if self.validate_headers_with_llm and candidates:
            validated_indices = await self._validate_headers_with_llm(
                candidates, slide_idx, split_elements, candidate_line_map
            )
        else:
            _, filtered_to_original, _ = self._pre_filter_header_candidates(candidates, candidate_line_map)
            validated_indices = set(filtered_to_original.values())

        result_elements = []
      

        for line_idx, line_elements in sorted(lines.items()):
            
            header_found = False
            for candidate_idx, info in candidate_line_map.items():
                if info['line_idx'] == line_idx and candidate_idx in validated_indices:
                    
                    leftmost_elem = info['line_elements'][0]

                    header_element = SlideElement(
                        element_type=ElementType.TEXTBOX,
                        bbox=BoundingBox(
                            x=leftmost_elem.bbox.x,
                            y=leftmost_elem.bbox.y,
                            width=sum(e.bbox.width for e in info['line_elements']),
                            height=max(e.bbox.height for e in info['line_elements'])
                        ),
                        content=info['combined_content'].strip(),
                        slide_index=leftmost_elem.slide_index,
                        z_index=leftmost_elem.z_index,
                        is_header=True,  
                    )

                    result_elements.append(header_element)  
                    header_found = True
                    remaining_elements = info['line_elements'][1:]
                    # print(f"[DEBUG] Remaining elements {remaining_elements}")
                    result_elements.extend(remaining_elements)
                    break

            if not header_found:

                result_elements.extend(line_elements)
            # print(result_elements)
            # print("\n")

        
        result_elements.sort(key=lambda e: (e.bbox.y, e.bbox.x))

        return result_elements

    async def _validate_headers_with_llm(
        self, candidates: List[str], slide_idx: int, slide_elements: List[SlideElement] = None,
        candidate_line_map: Dict[int, Dict] = None
    ) -> set:
        """
        Validate header candidates using LLM with minimal context.
        
        Optimizations:
        1. Pre-filters obvious non-headers to reduce candidates sent to LLM
        2. Uses position context only (TOP/BODY) instead of full slide structure

        Args:
            candidates: List of candidate header texts
            slide_idx: Slide index for context
            slide_elements: All elements from the slide
            candidate_line_map: Map of candidate index -> metadata including position info

        Returns:
            Set of validated candidate indices
        """
        if not candidates:
            return set()

        if self._llm_client is None:
            self._llm_client = LLMProvider.create()
        
        

        # Pre-filter obvious non-headers
        filtered_candidates, filtered_to_original, auto_rejected = self._pre_filter_header_candidates(
            candidates, candidate_line_map
        )
        
        # If all candidates were pre-filtered out, return empty set
        if not filtered_candidates:
            return set()
        
        # Build minimal context with position only
        candidate_lines = []
        for filtered_idx, (text, position) in enumerate(filtered_candidates):
            candidate_lines.append(f"{filtered_idx}. [{position}] \"{text}\"")

        candidates_text = "\n".join(candidate_lines)
        print("candidates_text:\n", candidates_text)
        system_prompt = f"""
        You are an expert at identifying Japanese/English document section headers.
        Determine if each candidate is a section HEADER.

        ## Candidates (with position: TOP=slide title area, BODY=content area)
        - *Atleast and only one candidate from [TOP] is allowed to be a header*, with priotize by the order of candidates.
        - For [BODY] candidates, multiple candidates can be a header.

        {candidates_text}

        ## HEADER Definition
        A HEADER is a CONCISE section title with number (may has many levels) + topic. Examples:
        - "4-1.2 User Settings" → YES
        - "5.1.6.6.1 対象地域　TARGET REGIONS" → YES
        - "・8. アーキテクチャ" → YES
        - "12.1. Common ▲h" → YES
        Sometimes PPTX files have abnormal header format with no numbered part at first, but they usually in "TOP" position

        ## NOT a HEADER
        - Excessive number heading (e.g, '526', '234') → NO
        - Table of contents, change history, version notes, etc. → NO
        - Long paragraphs describing PROCESSES or ACTIONS → NO
        - Steps/instructions/List items (e.g., "1. Some text, 2. Some text (2)") → NO
        - TOC entries, measurements, UI labels → NO

        ## CRITICAL: Sequential List Detection
        - There is no way in a slide, we can have so many headers. They most likely a list.
        - If you see 3+ candidates with sequential numbers (1., 2., 3.), mark ALL as NO.

        ## Output
        One YES or NO per line, in order. No explanations."""

        user_prompt = "Provide the 'YES', 'NO' result for each heading candidate."
        try:
            provider = os.getenv("LLM_PROVIDER", "openai").lower()
            if provider == "openai":
                model_config = {
                    "model": "gpt-5-mini",
                    "reasoning": {"effort": "low"},
                    "text": {"verbosity": "low"}
                }
            elif provider == "azure":
                model_config = {
                    "model": "gpt-5-mini",
                    "reasoning": {"effort": "low"},
                    "text": {"verbosity": "low"}
                }
            else:
                model_config = {"model": "apac.amazon.nova-pro-v1:0", "region": "ap-southeast-1"}

            response = await self._llm_client.create(
                **model_config,
                instructions=system_prompt,
                input=[{"role": "user", "content": [{"type": "input_text", "text": user_prompt}]}],
            )
            validated = set()
            response_lines = [line.strip().upper() for line in response.strip().split('\n') if line.strip()]
            for filtered_idx in range(len(filtered_candidates)):
                if filtered_idx < len(response_lines):
                    answer = response_lines[filtered_idx]
                    print("answer: ", answer)
                    if 'YES' in answer:
                        # Map back to original index
                        original_idx = filtered_to_original[filtered_idx]
                        validated.add(original_idx)

            return validated

        except Exception as e:
            print(f"Warning: LLM validation failed for slide {slide_idx + 1}: {e}")
            # On failure, validate all filtered candidates
            return set(filtered_to_original.values())

    def _pre_filter_header_candidates(
        self, 
        candidates: List[str], 
        candidate_line_map: Dict[int, Dict] = None
    ) -> Tuple[List[Tuple[str, str]], Dict[int, int], set]:
        """
        Pre-filter obvious non-headers before sending to LLM.
        
        Filters out:
        - Very long text (>200 chars) - likely paragraph content
        - Text with multiple sentences
        - Measurement/value patterns
        - Empty or whitespace-only text
        
        Args:
            candidates: List of candidate header texts
            candidate_line_map: Map of candidate index -> metadata
            
        Returns:
            Tuple of:
            - filtered_candidates: List of (text, position) tuples
            - filtered_to_original: Map of filtered index -> original index
            - auto_rejected: Set of original indices that were auto-rejected
        """
        candidates = [normalize_line(candidate) for candidate in candidates]
        filtered_candidates = []  # List of (text, position)
        filtered_to_original = {}  # filtered_idx -> original_idx
        auto_rejected = set()
        
        for original_idx, text in enumerate(candidates):
            text = text.strip() if text else ""
            
            # Get position info
            position = "BODY"
            if candidate_line_map and original_idx in candidate_line_map:
                info = candidate_line_map[original_idx]
                position = "TOP" if info.get('is_first_line', False) else "BODY"
            
            # Filter 1: Empty or too short
            if not text or len(text) < 2:
                auto_rejected.add(original_idx)
                continue
            
            # Filter 2: Very long text (likely paragraph content)
            if len(text) > 200:
                auto_rejected.add(original_idx)
                continue
            
            # Filter 3: Multiple sentences (period + space + capital or Japanese)
            # Pattern: ends sentence, followed by new sentence start
            sentence_breaks = re.findall(r'[.。]\s+[A-Z\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF]', text)
            if len(sentence_breaks) >= 2:
                auto_rejected.add(original_idx)
                continue
            
            # Filter 4: Pure measurement/value patterns
            if re.match(r'^[\d.,]+\s*(MHz|dB|km|kg|m|s|ms|%|℃|°C)\s*$', text, re.IGNORECASE):
                auto_rejected.add(original_idx)
                continue
            
            # Filter 5: Looks like a table data row (pipe-separated values)
            if text.count('|') >= 2:
                auto_rejected.add(original_idx)
                continue
            
            # PPTX files sometime have abnormal header format with no numbered part at first
            # Filter 6: Has no numbered part at first
            # if not re.match(r'^\d+(?:[.\-]\d+)*\.?(?:\s|$)', text):
            #     auto_rejected.add(original_idx)
            #     continue

            # Passed pre-filtering
            filtered_idx = len(filtered_candidates)
            filtered_candidates.append((text, position))
            filtered_to_original[filtered_idx] = original_idx

        return filtered_candidates, filtered_to_original, auto_rejected
        
    # =========================================================================
    # STEP 4.5: TOC-Based Tree Building (alternative to LLM heading detection)
    # =========================================================================

    def _build_tree_from_toc(
        self,
        toc_entries: List[Dict[str, Any]],
        all_elements: List[SlideElement],
    ) -> List[ContentNode]:
        """
        Build hierarchical tree from TOC entries by matching heading text
        against slide content. Unlike PDF, PPTX TOC entries lack page numbers,
        so we use sequential text matching.

        Algorithm:
        1. Build tree skeleton from toc_entries (stack-based, same as PDF reader)
        2. Sequential scan through all_elements:
           - Normalize text and compare against pending TOC headings
           - Switch current_node when a heading match is found
           - Non-matching content is appended to current_node
        3. Disambiguation: forward-only, standalone-line, top-area preference

        Args:
            toc_entries: List of dicts with heading_text, heading_level
            all_elements: All slide elements in order (without header marking)

        Returns:
            List of root ContentNode objects forming the tree
        """
        if not toc_entries:
            return []

        # --- Step 1: Build tree skeleton ---
        tree: List[ContentNode] = []
        stack: List[Tuple[int, ContentNode]] = []  # (level, node)
        node_list: List[ContentNode] = []  # Flat list for content assignment

        for entry in toc_entries:
            if entry.get("is_pre_chapter", False):
                print(f"Skipping pre-chapter entry: {entry.get('heading_text', '')}")
                continue

            level = entry["heading_level"]
            heading_text = entry["heading_text"]

            node = ContentNode(heading_text=heading_text, heading_level=level, content="")
            node_list.append(node)

            # Pop stack until we find a parent with lower level
            while stack and stack[-1][0] >= level:
                stack.pop()

            if stack:
                stack[-1][1].children.append(node)
            else:
                tree.append(node)

            stack.append((level, node))

        # --- Step 2: Build normalized heading map for matching ---
        def normalize_for_match(s: str) -> str:
            s = normalize_line(s)
            return re.sub(r'\s+', '', s).lower()

        normalized_headings: List[str] = [
            normalize_for_match(node.heading_text) for node in node_list
        ]

        # Track which TOC entry we're looking for next (forward-only)
        next_entry_idx = 0
        current_node: Optional[ContentNode] = None

        # --- Step 2b: Extract numbering prefixes for fallback matching ---
        # e.g. "4-1.2.32" from "4-1.2.32\t交通ラジオ機能（JP)"
        _NUM_PREFIX_RE = re.compile(
            r'^(\d+(?:[-．.‐－ ]\d+)+)'
        )

        def _extract_num_prefix(text: str) -> Optional[str]:
            """Extract numbering prefix like '4-1.2.32' from text."""
            cleaned = text.replace('－', '-').replace(' ', '-').replace('‐', '-').replace('．', '.')
            cleaned = re.sub(r'[\s\t　]+', '', cleaned)  # remove all whitespace
            m = _NUM_PREFIX_RE.match(cleaned)
            return m.group(1) if m else None

        heading_num_prefixes: List[Optional[str]] = [
            _extract_num_prefix(node.heading_text) for node in node_list
        ]

        # --- Step 2c: Detect slide-marker lines (appear on many slides) ---
        # Lines like "機  能  名  称", "Function Name", "備考", "Remark"
        # appear as template fields on every slide and should be skipped.
        line_freq: Dict[str, int] = {}
        total_element_count = 0
        for element in all_elements:
            content = element.content.strip()
            if not content:
                continue
            if (content.startswith(IMAGE_PATH_START_MARKER)
                    or content.startswith(TABLE_CONTENT_START_MARKER)
                    or content.startswith(IMAGE_CONTENT_START_MARKER)):
                continue
            total_element_count += 1
            for line in content.split('\n'):
                norm = normalize_for_match(line.strip())
                if norm and len(norm) < 20:
                    line_freq[norm] = line_freq.get(norm, 0) + 1

        marker_threshold = max(5, int(total_element_count * 0.10))
        marker_lines: set = {
            norm for norm, count in line_freq.items()
            if count >= marker_threshold
        }

        # --- Step 3: Sequential scan through all elements, LINE BY LINE ---
        # A single text shape may contain heading + sub-headings + content
        # all in one block. By processing line by line, each line gets its
        # own chance to match a pending TOC heading.

        def _append_to_node(node: ContentNode, text: str):
            """Helper to append text to a node's content."""
            if node.content:
                node.content += "\n" + text
            else:
                node.content = text

        def _try_match_heading(line_text: str) -> bool:
            """Try to match a text line against pending TOC headings.
            Returns True if matched, False otherwise.
            Updates nonlocal next_entry_idx and current_node."""
            nonlocal next_entry_idx, current_node

            norm_line = normalize_for_match(line_text)
            if not norm_line:
                return False

            # Extract numbering prefix from this line for fallback matching
            line_num_prefix = _extract_num_prefix(line_text)

            def _check_match(heading_idx: int) -> bool:
                norm_heading = normalized_headings[heading_idx]
                if not norm_heading or len(norm_heading) < 2:
                    return False
                # 1-3. Text-based match
                if (
                    norm_line == norm_heading
                    or (norm_line in norm_heading and len(norm_line) > 3)
                    or (norm_line.startswith(norm_heading) and len(norm_heading) > 3
                        and len(norm_line) <= len(norm_heading) * 1.5)
                ):
                    return True
                # 4. Numbering-prefix fallback
                if (
                    line_num_prefix
                    and heading_num_prefixes[heading_idx]
                    and line_num_prefix == heading_num_prefixes[heading_idx]
                    and len(line_num_prefix) >= 3
                ):
                    return True
                return False

            # Check if this line is just repeating the current heading
            # (e.g. at the top of a continuation slide). If so, skip it.
            if next_entry_idx > 0 and _check_match(next_entry_idx - 1):
                return True

            # Check if this line matches any pending heading
            for i in range(next_entry_idx, len(node_list)):
                if _check_match(i):
                    current_node = node_list[i]
                    next_entry_idx = i + 1
                    return True

            return False

        for element in all_elements:
            content = element.content.strip()
            if not content:
                continue

            # Non-text content (images, tables)   append as-is to current node
            is_non_text = (
                content.startswith(IMAGE_PATH_START_MARKER)
                or content.startswith(TABLE_CONTENT_START_MARKER)
                or content.startswith(IMAGE_CONTENT_START_MARKER)
            )

            if is_non_text:
                if current_node:
                    _append_to_node(current_node, content)
                continue

            # Split element content into individual lines for matching.
            # This handles the case where heading + sub-headings + body
            # are all inside a single text shape.
            lines = content.split('\n')

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Non-text markers within multi-line content
                if (
                    line.startswith(IMAGE_PATH_START_MARKER)
                    or line.startswith(TABLE_CONTENT_START_MARKER)
                    or line.startswith(IMAGE_CONTENT_START_MARKER)
                ):
                    if current_node:
                        _append_to_node(current_node, line)
                    continue

                # Skip slide-template marker lines (e.g. "機能名称",
                # "Function Name", "備考", "Remark")
                norm_line_check = normalize_for_match(line)
                if norm_line_check in marker_lines:
                    continue

                # Try to match this line against pending TOC headings
                if _try_match_heading(line):
                    continue

                # No heading match   append as content to current node
                if current_node:
                    _append_to_node(current_node, line)

        return tree

    def _build_candidate_frequency_map(
        self, prs: Presentation, total_slides: int
    ) -> Tuple[Dict[str, int], int]:
        """
        Pre-scan the first N slides to build a frequency map of top-area text.

        This detects repeated "marker" lines (company names, breadcrumbs,
        logos/labels) that appear at the top of many slides and would
        otherwise be incorrectly picked as headers by the LLM.

        Args:
            prs: Presentation object
            total_slides: Total number of slides

        Returns:
            Tuple of (frequency_map, slides_scanned)
            frequency_map: {normalized_text: count}
        """
        scan_count = max(5, int(total_slides * 0.20))
        scan_count = min(scan_count, total_slides)

        freq_map: Dict[str, int] = {}
        slides_scanned = 0

        for slide_idx, slide in enumerate(prs.slides):
            if slides_scanned >= scan_count:
                break
            # Skip TOC slides
            if self._is_toc_slide(slide_idx):
                continue

            slide_text_lines = self._extract_text_from_slide(slide)
            if not slide_text_lines:
                continue

            slides_scanned += 1
            # Collect candidates from top area (first line + top 20%)
            seen_on_this_slide: set = set()
            for line_idx, line in enumerate(slide_text_lines):
                line = line.strip()
                if not line:
                    continue
                # Only consider first line and first ~20% of lines
                is_first_line = (line_idx == 0)
                is_top_area = (line_idx < max(1, len(slide_text_lines) * 0.20))

                if is_first_line or is_top_area:
                    norm = normalize_line(line).strip().lower()
                    if norm and norm not in seen_on_this_slide:
                        seen_on_this_slide.add(norm)
                        freq_map[norm] = freq_map.get(norm, 0) + 1

        return freq_map, slides_scanned

    # =========================================================================
    # STEP 5: Build Tree Structure
    # =========================================================================

    def _extract_header_level(self, header_text: str) -> int:
        """
        Determine header level from the header text pattern.
        Uses the heading_detector module for consistent detection.
        """
        header_text = (header_text or "").strip()
        if not header_text:
            return 1

        # Use heading_detector module for level detection
        level, pattern_name, match = detect_heading_level(header_text)
        
        if level is not None:
            return level
        
        # Fallback heuristic: count dots in leading numeric outline
        # Normalize the text first to handle full-width characters
        normalized = normalize_line(header_text)
        m = re.match(r'^(\d+(?:[.\-]\d+)*)', normalized)
        if m:
            return m.group(1).count(".") + 1

        return 1

    def _build_tree_from_elements(
        self,
        elements: List[SlideElement],
    ) -> List[ContentNode]:
        """
        Build hierarchical tree from elements with validated headers.

        Elements are already processed with headers split and marked with is_header=True.
        Simply traverse elements in order and build tree based on header property and levels.

        Args:
            elements: List of slide elements (headers already split and validated)

        Returns:
            List of root ContentNode objects forming the tree
        """
        if not elements:
            return []

        nodes = []
        current_node: Optional[ContentNode] = None
        node_stack: List[ContentNode] = [] 

        def add_node(header_text: str, level: int, previous_content: str):
            """Helper to add a header node to the tree."""
            nonlocal current_node, node_stack, nodes
            new_node = ContentNode(
                heading_text=header_text,
                heading_level=level,
                content=previous_content,
            )
            while node_stack and node_stack[-1].heading_level >= level:
                node_stack.pop()

            if node_stack:
                node_stack[-1].children.append(new_node)
            else:
                nodes.append(new_node)

            node_stack.append(new_node)
            current_node = new_node

        def append_content(content: str):
            """Helper to append content to current node."""
            nonlocal current_node, node_stack, nodes
            
            if not content:
                return

            if current_node:
                if current_node.content:
                    current_node.content += "\n" + content
                else:
                    current_node.content = content
                    
            # comment this to avoid orphan node
            
            # else:
            #     current_node = ContentNode(
            #         heading_text="",
            #         heading_level=1,
            #         content=content,
            #     )
            #     nodes.append(current_node)
            #     node_stack.append(current_node)
        header_flag = False
        reserved_contents =[]
        for element in elements:
            content = element.content.strip()
            if not content:
                continue
            if element.is_header:
                header_flag = True
                previous_content = ""
                for reserved_content in reserved_contents:
                    previous_content += reserved_content
                reserved_contents = []
                level = self._extract_header_level(content)
                
                # treat header level > 3 as normal content to avoid deep nesting
                if level > 3:
                    content = previous_content + content
                    append_content(content)
                else:
                    add_node(content, level, previous_content)
                
                    
            else:
                if not header_flag:
                    reserved_contents.append(content)
                append_content(content)
        return nodes

    # =========================================================================
    # Main Entry Point
    # =========================================================================

    @staticmethod
    def _detect_and_convert_ppt(path: Path, temp_dir: Path) -> Path:
        """Convert legacy .ppt format to .pptx using LibreOffice.

        Returns the path to the converted .pptx file inside *temp_dir*.
        The caller is responsible for cleaning up *temp_dir*.
        """
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx",
             "--outdir", str(temp_dir), str(path)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice .ppt→.pptx conversion failed: {result.stderr}")
        converted = temp_dir / (path.stem + ".pptx")
        if not converted.exists():
            raise FileNotFoundError(f"Converted file not found: {converted}")
        return converted

    def extract_pptx_text(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Extract text and content from a PPTX (or legacy .ppt) file.

        If the input is a .ppt file, it is first converted to .pptx via
        LibreOffice, then parsed normally.

        Dual-path approach:
        A) If TOC slides are detected and TOC entries extracted:
           1. Build slide structure, merge elements (no LLM heading detection)
           2. Build tree from TOC entries by text-matching against slide content
        B) If no TOC found (fallback):
           1. Pre-scan slides for repeated marker lines in top area
           2. Build slide structure, merge elements
           3. Detect headers with LLM (filtered by frequency map)
           4. Build tree from validated headers

        Args:
            pptx_path: Path to the PowerPoint file

        Returns:
            List of dictionaries representing the document tree structure
        """
        original_path = Path(pptx_path)
        temp_dir_obj = None

        # Legacy .ppt → convert to .pptx first
        if original_path.suffix.lower() == ".ppt":
            temp_dir_obj = Path(tempfile.mkdtemp(prefix="ppt2pptx_"))
            try:
                converted = self._detect_and_convert_ppt(original_path, temp_dir_obj)
                pptx_path = str(converted)
            except Exception:
                shutil.rmtree(temp_dir_obj, ignore_errors=True)
                raise

        try:
            return self._extract_pptx_text_inner(pptx_path)
        finally:
            if temp_dir_obj is not None:
                shutil.rmtree(temp_dir_obj, ignore_errors=True)

    def _extract_pptx_text_inner(self, pptx_path: str) -> List[Dict[str, Any]]:
        """Core extraction logic for a .pptx file (already in OOXML format)."""
        import asyncio

        self._pptx_path = os.path.abspath(pptx_path)

        prs = Presentation(pptx_path)

        self._slide_width = prs.slide_width or 0
        self._slide_height = prs.slide_height or 0
        file_name = os.path.basename(pptx_path)
        print(f"[PPTX Parser] Start parsing {file_name}")

        # Step 0: Detect TOC slides AND extract TOC entries
        self._toc_slides, toc_entries = self._detect_toc_slides(prs)
        use_toc_strategy = len(toc_entries) > 0

        if use_toc_strategy:
            print(f"[PPTX Parser] Using TOC-based heading tree ({len(toc_entries)} entries)")
        else:
            print("[PPTX Parser] No TOC entries found, using LLM-based heading detection")

        # For LLM fallback: pre-scan for repeated marker lines
        self._marker_freq_map: Dict[str, int] = {}
        self._marker_scan_count: int = 0
        if not use_toc_strategy:
            total_slides = len(prs.slides)
            self._marker_freq_map, self._marker_scan_count = (
                self._build_candidate_frequency_map(prs, total_slides)
            )
            if self._marker_freq_map:
                print(f"[PPTX Parser] Pre-scanned {self._marker_scan_count} slides for repeated markers, "
                      f"found {len(self._marker_freq_map)} unique top-area texts")

        all_elements = []

        try:
            loop = asyncio.get_event_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)

        start_time = time.time()
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_master = prs.slide_masters[0]
            if slide_idx % 10 == 0 or slide_idx == len(prs.slides):
                print(f"[PPTX Parser] Successfully processed {slide_idx} slides")
            # Skip TOC slides
            if self._is_toc_slide(slide_idx - 1):
                print(f"[TOC] Skipping TOC slide {slide_idx}")
                continue
            # Step 1: build slide structure
            elements = self._build_slide_structure(slide, slide_idx, slide_master)
            # Step 2: merge elements
            print(f"[DEBUG] Working for merged element for slide index : {slide_idx}")
            elements = self._merge_elements(elements)
            # print(f"[DEBUG] Elements before splitting headers : {elements}")
            if not use_toc_strategy:
                # Step 3 (fallback only): detect and split headers with LLM validation
                try:
                    print(f"[DEBUG] Before detect and split headers in slide : {len(elements)}")
                    elements = loop.run_until_complete(
                        self._detect_and_split_headers_in_slide(elements, slide_idx)
                    )                    

                except Exception as e:
                    print(f"Warning: Header detection failed for slide {slide_idx + 1}: {e}")
                    import traceback
                    traceback.print_exc()

            all_elements.extend(elements)
        end_time = time.time()
        print(f"[PPTX Parser] Successfully processed {file_name} in {end_time - start_time:.2f} seconds")

        # Step 4: Merge continuation tables (tables spanning multiple slides)
        all_elements = self._merge_continuation_tables(all_elements)

        
        # Step 5: Build tree structure
        if use_toc_strategy:
            print("Using TOC")  
            tree = self._build_tree_from_toc(toc_entries, all_elements)
        else:
            print("Not using TOC")
            tree = self._build_tree_from_elements(all_elements)
        # print(f"[DEBUG] Tree : {tree}")
        result = [node.to_dict() for node in tree]
        # print(f"[DEBUG] Result: {result}")
        if not result:
            all_parts = [elem.content for elem in all_elements if elem.content]
            result = [{
                "heading_text": "",
                "heading_level": 0,
                "content": "\n".join(all_parts),
                "children": []
            }]
        return result

    async def extract_pptx_text_async(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Async version of extract_pptx_text.

        Args:
            pptx_path: Path to the PowerPoint file

        Returns:
            List of dictionaries representing the document tree structure
        """
        self._pptx_path = os.path.abspath(pptx_path)

        prs = Presentation(pptx_path)
        self._slide_width = prs.slide_width or 0
        self._slide_height = prs.slide_height or 0
        file_name = os.path.basename(pptx_path)
        print(f"[PPTX Parser] Start parsing {file_name}")

        # Step 0: Detect TOC slides AND extract TOC entries
        self._toc_slides, toc_entries = self._detect_toc_slides(prs)
        use_toc_strategy = len(toc_entries) > 0

        if use_toc_strategy:
            print(f"[PPTX Parser] Using TOC-based heading tree ({len(toc_entries)} entries)")
        else:
            print("[PPTX Parser] No TOC entries found, using LLM-based heading detection")

        # For LLM fallback: pre-scan for repeated marker lines
        self._marker_freq_map: Dict[str, int] = {}
        self._marker_scan_count: int = 0
        if not use_toc_strategy:
            total_slides = len(prs.slides)
            self._marker_freq_map, self._marker_scan_count = (
                self._build_candidate_frequency_map(prs, total_slides)
            )

        all_elements = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            if slide_idx % 10 == 0 or slide_idx == len(prs.slides):
                print(f"[PPTX Parser] Successfully processed {slide_idx} slides")
            if self._is_toc_slide(slide_idx - 1):
                print(f"[TOC] Skipping TOC slide {slide_idx}")
                continue
            # Step 1: Build slide structure
            elements = self._build_slide_structure(slide, slide_idx)
            # Step 2: Merge elements
            elements = self._merge_elements(elements)

            if not use_toc_strategy:
                # Step 3 (fallback only): Detect and split headers with LLM validation
                elements = await self._detect_and_split_headers_in_slide(elements, slide_idx)

            all_elements.extend(elements)
        
        # Step 4: Merge continuation tables (tables spanning multiple slides)
        all_elements = self._merge_continuation_tables(all_elements)
        
        # Step 5: Build tree structure
        if use_toc_strategy:
            tree = self._build_tree_from_toc(toc_entries, all_elements)
        else:
            tree = self._build_tree_from_elements(all_elements)

        return [node.to_dict() for node in tree]
        
    def cleanup(self):
        """Clean up temporary files."""
        if self._temp_dir and os.path.exists(self._temp_dir):
            try:
                shutil.rmtree(self._temp_dir)
            except Exception as e:
                print(f"Warning: Failed to cleanup temp dir: {e}")
        if self._slide_exports_dir and os.path.exists(self._slide_exports_dir):
            try:
                shutil.rmtree(self._slide_exports_dir)
            except Exception as e:
                print(f"Warning: Failed to cleanup slide exports dir: {e}")
        self._temp_dir = None
        self._slide_exports_dir = None
        self._slide_export_cache = {}
        self._saved_images = []
        self._image_counter = 0
        self._pptx_path = None

    def __del__(self):
        """Destructor to ensure cleanup (unless keep_temp_files is True)."""
        if not getattr(self, 'keep_temp_files', False):
            self.cleanup()

# for unit text only
def extract_pptx_text(
    pptx_path: str,
    get_table_content: bool = True,
    get_image_content: bool = True,
    get_shape_content: bool = True,
    validate_headers_with_llm: bool = True,
    keep_temp_files: bool = True,
    skip_toc: bool = True,
) -> List[Dict[str, Any]]:
    parser = PptxParser(
        get_table_content=get_table_content,
        get_image_content=get_image_content,
        get_shape_content=get_shape_content,
        validate_headers_with_llm=validate_headers_with_llm,
        keep_temp_files=keep_temp_files,
        skip_toc=skip_toc,
        repository_id='',
    )
    return parser.extract_pptx_text(pptx_path)
