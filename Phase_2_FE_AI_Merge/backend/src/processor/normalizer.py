"""
Document Normalizer Module

This module normalizes various document formats into either PDF or Markdown format
to prepare them for the document processing pipeline.

Supported Input Formats:
- Documents: DOCX, PPTX, ODT, RTF
- Web: HTML, MHTML
- Spreadsheets: XLSX, CSV
- Images: PNG, JPG, JPEG, BMP, TIFF
- Already normalized: PDF, MD, TXT

Output Formats:
- PDF: For image-based RAG retrieval (via Docling OCR)
- Markdown/Text: For text-based RAG retrieval
"""

import os
import json
import shutil
import base64
from pathlib import Path
from typing import Dict, List, Optional, Union, Tuple
from dataclasses import dataclass
import logging
from datetime import datetime
from tqdm import tqdm
import tempfile

from src.processor.utils import sanitize_filename_stem

# Import conversion libraries
try:
    from docx import Document
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import img2pdf
    IMG2PDF_AVAILABLE = True
except ImportError:
    IMG2PDF_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    from reportlab.lib import colors
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False


@dataclass
class NormalizerConfig:
    """Configuration for document normalization."""
    
    # Output format preferences
    generate_pdf: bool = True          # Generate PDF for image-based retrieval
    generate_markdown: bool = True     # Generate markdown for text-based retrieval
    
    # Image handling
    image_to_pdf: bool = True          # Convert standalone images to PDF
    image_quality: int = 95            # JPEG quality for image conversion
    
    # Spreadsheet handling
    excel_to_markdown: bool = True     # Convert Excel to Markdown tables
    csv_to_markdown: bool = True       # Convert CSV to Markdown tables
    max_table_rows: int = 1000         # Maximum rows to include in markdown
    
    # HTML handling
    html_preserve_formatting: bool = True  # Keep some HTML formatting in markdown
    
    # Document handling
    docx_extract_images: bool = True   # Extract embedded images from DOCX
    pptx_extract_images: bool = True   # Extract images from slides
    
    # PDF generation settings
    pdf_page_size: str = "A4"          # A4 or letter
    pdf_margin: float = 0.75           # Margin in inches
    pdf_content_source: str = "hybrid" # pymupdf | docling | hybrid for born-digital PDF content
    pdf_reader_enable_ocr: bool = False
    pdf_reader_extract_images: bool = False
    runtime_yaml: Optional[Dict] = None
    
    # Output organization
    organize_by_type: bool = True      # Organize output by original file type


class DocumentNormalizer:
    """
    Normalizes various document formats into PDF and/or Markdown for RAG pipeline.
    
    This is the first stage in the processing pipeline:
    1. Normalizer converts everything to PDF/Markdown
    2. Document Processor processes the normalized files
    """
    
    def __init__(
        self,
        input_dir: Union[str, Path],
        output_dir: Union[str, Path],
        config: Optional[NormalizerConfig] = None
    ):
        """
        Initialize the Document Normalizer.
        
        Args:
            input_dir: Directory containing source documents
            output_dir: Directory for normalized outputs
            config: Normalization configuration
        """
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self.config = config or NormalizerConfig()
        
        # Create output directories
        self.pdf_dir = self.output_dir / "normalized_pdfs"
        self.markdown_dir = self.output_dir / "normalized_markdown"
        self.metadata_dir = self.output_dir / "normalization_metadata"
        self.originals_dir = self.output_dir / "original_files"  # Keep originals for Docling Stage 3
        self.excel_parsed_dir = self.output_dir / "excel_parsed"  # Custom Excel parser output
        self.docx_parsed_dir = self.output_dir / "docx_parsed"   # Custom DOCX parser output
        self.pdf_parsed_dir = self.output_dir / "pdf_parsed"     # Custom PDF parser output (born-digital)

        for dir_path in [self.pdf_dir, self.markdown_dir, self.metadata_dir, self.originals_dir, self.excel_parsed_dir, self.docx_parsed_dir, self.pdf_parsed_dir]:
            dir_path.mkdir(parents=True, exist_ok=True)
        
        # Setup logging
        self._setup_logging()
        
        # Track statistics
        self.stats = {
            "total_files": 0,
            "normalized_files": 0,
            "failed_files": 0,
            "by_type": {},
            "errors": []
        }
    
    def _setup_logging(self):
        """Setup logging configuration."""
        log_file = self.output_dir / "normalization.log"
        
    def normalize_batch(self) -> Dict:
        """
        Normalize all files in the input directory.
        
        Returns:
            Dictionary with processing statistics
        """
        print(f"Starting batch normalization from: {self.input_dir}")
        print(f"Output directory: {self.output_dir}")
        
        # Find all files to process
        files_to_process = self._find_files()
        self.stats["total_files"] = len(files_to_process)
        
        print(f"Found {len(files_to_process)} files to normalize")
        
        # Process each file
        for file_path in tqdm(files_to_process, desc="Normalizing documents"):
            try:
                self._normalize_file(file_path)
                self.stats["normalized_files"] += 1
            except Exception as e:
                print(f"ERROR: Failed to normalize {file_path}: {str(e)}")
                self.stats["failed_files"] += 1
                self.stats["errors"].append({
                    "file": str(file_path),
                    "error": str(e)
                })
        
        # Save statistics
        self._save_statistics()
        
        print(f"Normalization complete. {self.stats['normalized_files']}/{self.stats['total_files']} files processed successfully")
        
        return self.stats
    
    def _find_files(self) -> List[Path]:
        """Find all supported files in input directory."""
        supported_extensions = {
            # Documents
            '.docx', '.doc', '.odt', '.rtf',
            # Presentations
            '.pptx', '.ppt', '.odp',
            # Web
            '.html', '.htm', '.mhtml', '.xhtml',  # Added .xhtml
            # Spreadsheets
            '.xlsx', '.xls', '.xlsm', '.csv',
            # Images
            '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp',  # Added .webp
            # Already normalized (just copy)
            '.pdf', '.md', '.txt',
            # Text formats
            '.adoc', '.asciidoc', '.asc',  # Added AsciiDoc
            '.vtt'  # Added WebVTT
        }
        
        files = []
        for ext in supported_extensions:
            files.extend(self.input_dir.rglob(f"*{ext}"))
        
        return sorted(files)
    
    def _normalize_file(self, file_path: Path):
        """
        Normalize a single file based on its type.
        
        Args:
            file_path: Path to the file to normalize
        """
        ext = file_path.suffix.lower()
        stem = file_path.stem
        
        # Track by type
        file_type = self._get_file_type(ext)
        self.stats["by_type"][file_type] = self.stats["by_type"].get(file_type, 0) + 1
        
        print(f"Normalizing {file_type}: {file_path.name}")
        
        # Get safe filename (truncate if too long) - Windows 260 char path limit
        safe_stem = self._get_safe_filename(stem)
        
        # Copy ALL files except media to original_files for Stage 3 Docling processing
        # Use safe_stem to avoid Windows path length issues
        media_extensions = {'.mp4', '.avi', '.mov', '.mp3', '.wav', '.aac', '.ogg', '.flac', '.m4a'}
        if ext not in media_extensions:
            original_copy = self.originals_dir / f"{safe_stem}{ext}"
            shutil.copy2(file_path, original_copy)
            print(f"  → Saved original for Docling: {original_copy.name}")
        
        # Dispatch to appropriate handler with safe filename
        if ext in ['.docx', '.doc']:
            self._normalize_docx(file_path, safe_stem)
        elif ext in ['.pptx', '.ppt']:
            self._normalize_pptx(file_path, safe_stem)
        elif ext in ['.html', '.htm', '.mhtml', '.xhtml']:
            self._normalize_html(file_path, safe_stem)
        elif ext in ['.xlsx', '.xlsm', '.xls']:
            self._normalize_excel(file_path, safe_stem)
        elif ext == '.csv':
            self._copy_to_originals_only(file_path, safe_stem)  # CSV: only copy to originals
        elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp']:
            self._normalize_image(file_path, safe_stem)
        elif ext == '.pdf':
            self._copy_pdf(file_path, safe_stem)
        elif ext == '.txt':
            self._txt_to_md(file_path, safe_stem)
        elif ext == '.md':
            self._copy_text(file_path, safe_stem)
        elif ext in ['.adoc', '.asciidoc', '.asc', '.vtt']:
            self._copy_to_originals_only(file_path, safe_stem)  # AsciiDoc/WebVTT: only copy to originals
        else:
            print(f"WARNING: Unsupported file type: {ext}")
    
    def _get_file_type(self, ext: str) -> str:
        """Get file type category from extension."""
        type_map = {
            '.docx': 'document', '.doc': 'document', '.odt': 'document', '.rtf': 'document',
            '.pptx': 'presentation', '.ppt': 'presentation', '.odp': 'presentation',
            '.html': 'web', '.htm': 'web', '.mhtml': 'web', '.xhtml': 'web',
            '.xlsx': 'spreadsheet', '.xls': 'spreadsheet', '.xlsm': 'spreadsheet',
            '.csv': 'csv',
            '.png': 'image', '.jpg': 'image', '.jpeg': 'image', '.bmp': 'image',
            '.tiff': 'image', '.tif': 'image', '.webp': 'image',
            '.pdf': 'pdf',
            '.adoc': 'asciidoc', '.asciidoc': 'asciidoc', '.asc': 'asciidoc',
            '.vtt': 'webvtt',
            '.md': 'markdown', '.txt': 'text'
        }
        return type_map.get(ext.lower(), 'unknown')
    
    def _get_safe_filename(self, stem: str, max_length: int = 50) -> str:
        """
        Get a safe filename by truncating long names and adding hash for uniqueness.
        
        Args:
            stem: Original filename stem (without extension)
            max_length: Maximum length for filename (default 50 chars for Windows path limits)
            
        Returns:
            Safe filename stem, truncated if needed with MD5 hash appended
        """
        stem = sanitize_filename_stem(stem)
        if len(stem) <= max_length:
            return stem
        
        # Truncate and add hash for uniqueness
        import hashlib
        hash_suffix = hashlib.md5(stem.encode()).hexdigest()[:8]
        truncated = stem[:max_length - 9]  # Leave room for underscore and hash
        return f"{truncated}_{hash_suffix}"
    
    # ======================== DOCX Normalization ========================
    
    def _normalize_docx(self, file_path: Path, stem: str):
        """Normalize DOCX files with both parser and PDF conversion.

        Produces structured JSON in docx_parsed/ for heading-aware chunking,
        and also attempts PDF normalization via LibreOffice (same behavior
        pattern as PPTX) for visual preview parity.
        """
        try:
            from .docx_reader_v2 import DocxParser
            import json as _json

            parser = DocxParser(drop_deleted_table_content=True)
            tree = parser.extract_docx_text(str(file_path))

            out_json = self.docx_parsed_dir / f"{stem}.json"
            with open(out_json, "w", encoding="utf-8") as fh:
                _json.dump(tree, fh, ensure_ascii=False, indent=2)

            print(f"  → ✓ DOCX parsed to JSON: {out_json.name}")

            # Always attempt DOCX -> PDF normalization (parity with PPTX flow).
            if self.config.generate_pdf:
                print("  → Trying LibreOffice for DOCX to PDF conversion...")
                if not self._docx_to_pdf_libreoffice(file_path, stem):
                    # Fallback to ReportLab text-only PDF if python-docx is available.
                    if PYTHON_DOCX_AVAILABLE:
                        try:
                            doc = Document(file_path)
                            self._docx_to_pdf(doc, stem)
                        except Exception as e2:
                            print(f"  → ✗ DOCX PDF fallback failed: {e2}")
                    else:
                        print("WARNING: python-docx not available, DOCX PDF fallback skipped")
                else:
                    print("  → ✓ LibreOffice conversion successful (images preserved)")

        except Exception as e:
            print(f"  → ✗ Custom DOCX parser failed: {e}")
            print(f"  → Falling back to LibreOffice PDF conversion")
            if not PYTHON_DOCX_AVAILABLE:
                print("WARNING: python-docx not available, skipping DOCX fallback")
                return
            try:
                doc = Document(file_path)
                if self.config.generate_pdf:
                    if not self._docx_to_pdf_libreoffice(file_path, stem):
                        self._docx_to_pdf(doc, stem)
                    else:
                        print("  → ✓ LibreOffice conversion successful (images preserved)")
            except Exception as e2:
                print(f"ERROR: Error processing DOCX {file_path}: {str(e2)}")
                raise
    
    # Unused anymore but maybe keep
    def _docx_to_markdown(self, doc: 'Document') -> str:
        """Convert DOCX document to Markdown."""
        markdown_lines = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                markdown_lines.append("")
                continue
            
            # Check for heading style
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading ', '')
                try:
                    level = int(level)
                    markdown_lines.append(f"{'#' * level} {text}\n")
                except:
                    markdown_lines.append(f"{text}\n")
            else:
                markdown_lines.append(f"{text}\n")
        
        # Extract tables
        for table in doc.tables:
            markdown_lines.append("\n")
            markdown_lines.append(self._table_to_markdown(table))
            markdown_lines.append("\n")
        
        return "\n".join(markdown_lines)
    
    # Unused anymore but maybe keep
    def _table_to_markdown(self, table) -> str:
        """Convert a DOCX table to Markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        
        if len(rows) > 1:
            # Add separator after header
            header = rows[0]
            separator = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            return "\n".join([header, separator] + rows[1:])
        elif rows:
            return rows[0]
        return ""
    
    # Convert docx to pdf using libreoffice
    def _docx_to_pdf_libreoffice(self, file_path: Path, stem: str) -> bool:
        """Convert office-like files to PDF.

        Conversion order:
        1) AWS Lambda LibreOffice (when enabled)
        2) Local LibreOffice binary
        """
        # Use safe stem to avoid Windows path length issues
        pdf_path = self.pdf_dir / f"{stem}.pdf"

        print(f"    [LibreOffice] Converting {file_path.name} to PDF...")

        try:
            if self._should_use_lambda_office_pdf_conversion():
                print("    → Trying AWS Lambda LibreOffice conversion...")
                if self._office_to_pdf_lambda(file_path, pdf_path):
                    print("    ✓ AWS Lambda conversion successful")
                    return True
                print("    ✗ AWS Lambda conversion failed; trying local LibreOffice...")

            print("    → Trying local LibreOffice...")
            return self._office_to_pdf_local(file_path, pdf_path)

        except Exception as e:
            print(f"    ✗ LibreOffice error: {e}")
            return False

    def _should_use_lambda_office_pdf_conversion(self) -> bool:
        """Decide if Office->PDF should be delegated to Lambda."""
        mode = os.getenv("OFFICE_PDF_CONVERTER_MODE", "auto").strip().lower()
        if mode not in {"auto", "local", "lambda"}:
            mode = "auto"

        fn_name = os.getenv("OFFICE_PDF_LAMBDA_FUNCTION_NAME", "").strip()
        if not fn_name:
            print(f"    [Lambda] Disabled: OFFICE_PDF_LAMBDA_FUNCTION_NAME not set")
            return False

        if mode == "lambda":
            print(f"    [Lambda] FORCED mode: will use Lambda ({fn_name})")
            return True
        if mode == "local":
            print(f"    [Lambda] LOCAL mode: will use local LibreOffice")
            return False

        # Auto mode: use Lambda when running deployed cloud workload.
        storage_backend = os.getenv("FILE_STORAGE_BACKEND", "").strip().lower()
        in_cloud_runtime = bool(
            os.getenv("AWS_EXECUTION_ENV")
            or os.getenv("ECS_CONTAINER_METADATA_URI")
            or os.getenv("ECS_CONTAINER_METADATA_URI_V4")
        )
        use_lambda = storage_backend == "s3" or in_cloud_runtime
        print(f"    [Lambda] AUTO mode: storage={storage_backend}, in_cloud={in_cloud_runtime}, use_lambda={use_lambda}")
        return use_lambda

    def _office_to_pdf_lambda(self, file_path: Path, pdf_path: Path) -> bool:
        """Convert a file to PDF by invoking an AWS Lambda function.

        Expected Lambda response shape:
        {
          "ok": true,
          "pdf_base64": "<base64-encoded-pdf-bytes>"
        }
        """
        fn_name = os.getenv("OFFICE_PDF_LAMBDA_FUNCTION_NAME", "").strip()
        if not fn_name:
            print("    ✗ Lambda function name not configured")
            return False

        region = (
            os.getenv("OFFICE_PDF_LAMBDA_REGION", "").strip()
            or os.getenv("AWS_REGION", "").strip()
            or None
        )

        print(f"    [Lambda] Function: {fn_name}, Region: {region}")

        try:
            import boto3
        except Exception:
            print("    ✗ boto3 is not available; cannot invoke Lambda converter")
            return False

        try:
            file_bytes = file_path.read_bytes()
            file_size_kb = len(file_bytes) / 1024
            print(f"    [Lambda] Encoding file: {file_path.name} ({file_size_kb:.1f} KB)")

            payload = {
                "operation": "convert-to-pdf",
                "filename": file_path.name,
                "extension": file_path.suffix.lower(),
                "content_base64": base64.b64encode(file_bytes).decode("ascii"),
            }

            print(f"    [Lambda] Invoking {fn_name}...")
            client = boto3.client("lambda", region_name=region)
            response = client.invoke(
                FunctionName=fn_name,
                InvocationType="RequestResponse",
                Payload=json.dumps(payload).encode("utf-8"),
            )

            status_code = response.get("StatusCode", 0)
            print(f"    [Lambda] Response StatusCode: {status_code}")

            if int(status_code) != 200:
                print(f"    ✗ Lambda invoke failed, StatusCode={status_code}")
                return False

            raw = response.get("Payload").read() if response.get("Payload") else b"{}"
            body = json.loads(raw.decode("utf-8") or "{}")

            # Support API Gateway style payload from Lambda handlers.
            if isinstance(body, dict) and isinstance(body.get("body"), str):
                try:
                    body = json.loads(body["body"])
                except Exception:
                    pass

            if isinstance(body, dict) and body.get("ok") is False:
                error = body.get('error', 'unknown error')
                print(f"    ✗ Lambda converter returned error: {error}")
                return False

            pdf_b64 = None
            if isinstance(body, dict):
                pdf_b64 = body.get("pdf_base64") or body.get("pdfBase64")
                if not pdf_b64 and isinstance(body.get("data"), dict):
                    pdf_b64 = body["data"].get("pdf_base64") or body["data"].get("pdfBase64")

            if not pdf_b64:
                print("    ✗ Lambda converter response missing pdf_base64")
                print(f"    [Lambda] Response body keys: {list(body.keys()) if isinstance(body, dict) else 'not a dict'}")
                return False

            print(f"    [Lambda] Decoding PDF response...")
            pdf_bytes = base64.b64decode(pdf_b64)
            pdf_path.write_bytes(pdf_bytes)

            pdf_size_kb = len(pdf_bytes) / 1024
            print(f"    ✓ Lambda PDF created: {pdf_path.name} ({pdf_size_kb:.1f} KB)")

            return pdf_path.exists() and pdf_path.stat().st_size > 0

        except Exception as e:
            print(f"    ✗ Lambda converter error: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _office_to_pdf_local(self, file_path: Path, pdf_path: Path) -> bool:
        """Convert office-like files to PDF using a local LibreOffice binary."""
        try:
            import subprocess
            import sys

            if sys.platform == 'win32':
                soffice_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                ]
            elif sys.platform == 'darwin':
                soffice_paths = [
                    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                ]
            else:
                soffice_paths = [
                    "/usr/bin/soffice",
                    "/usr/local/bin/soffice",
                    "soffice",
                ]

            soffice = None
            for path in soffice_paths:
                if path == "soffice" or Path(path).exists():
                    soffice = path
                    break

            if not soffice:
                print("    ✗ LibreOffice not found in system")
                print("    ℹ Install LibreOffice from: https://www.libreoffice.org/download/")
                return False

            print(f"    ✓ Found LibreOffice: {soffice}")

            cmd = [
                soffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.pdf_dir),
                str(file_path)
            ]

            conversion_kwargs = {'capture_output': True, 'timeout': 60}

            if sys.platform == 'win32':
                startupinfo = subprocess.STARTUPINFO()
                startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                startupinfo.wShowWindow = subprocess.SW_HIDE
                conversion_kwargs['startupinfo'] = startupinfo
                conversion_kwargs['creationflags'] = subprocess.CREATE_NO_WINDOW

            subprocess.run(cmd, **conversion_kwargs)

            output_pdf = self.pdf_dir / f"{file_path.stem}.pdf"
            if output_pdf.exists() and output_pdf != pdf_path:
                output_pdf.rename(pdf_path)

            if pdf_path.exists():
                return True

            print("    ✗ LibreOffice conversion failed (output not found)")
            return False

        except Exception as e:
            print(f"    ✗ LibreOffice local conversion error: {e}")
            return False
    
    # Fallback PDF conversion using ReportLab (text-only)
    def _docx_to_pdf(self, doc: 'Document', stem: str):
        """Convert DOCX document to PDF using ReportLab (fallback only).
        Uses safe stem to avoid Windows path length issues.
        """
        if not REPORTLAB_AVAILABLE:
            print("WARNING: ReportLab not available, cannot convert DOCX to PDF")
            return
        
        # Use safe stem
        pdf_path = self.pdf_dir / f"{stem}.pdf"
        
        try:
            # Create PDF document
            pdf_doc = SimpleDocTemplate(
                str(pdf_path),
                pagesize=A4 if self.config.pdf_page_size == "A4" else letter,
                leftMargin=self.config.pdf_margin * inch,
                rightMargin=self.config.pdf_margin * inch,
                topMargin=self.config.pdf_margin * inch,
                bottomMargin=self.config.pdf_margin * inch
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            # Add paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Determine style based on heading
                    if para.style.name.startswith('Heading'):
                        style = styles['Heading1']
                    else:
                        style = styles['Normal']
                    
                    story.append(Paragraph(text, style))
                    story.append(Spacer(1, 0.2 * inch))
            
            pdf_doc.build(story)
            print(f"Created PDF: {pdf_path}")
        
        except Exception as e:
            print(f"ERROR: Error creating PDF from DOCX: {str(e)}")
    
    # ======================== PPTX Normalization ========================
    
    def _normalize_pptx(self, file_path: Path, stem: str):
        """Normalize PPTX files to PDF and Markdown."""
        if not PYTHON_PPTX_AVAILABLE:
            print("WARNING: python-pptx not available, skipping PPTX processing")
            return
        
        try:
            prs = Presentation(file_path)
            
            # Convert to PDF only - Try LibreOffice first (best quality, preserves images)
            if self.config.generate_pdf:
                print("  → Trying LibreOffice for PDF conversion...")
                if not self._docx_to_pdf_libreoffice(file_path, stem):
                    # Fallback to ReportLab if LibreOffice not available
                    print("  → LibreOffice not available, using ReportLab (text-only)")
                    self._pptx_to_pdf(prs, stem)
                else:
                    print("  → ✓ LibreOffice conversion successful (images preserved)")
                    
        except Exception as e:
            print(f"ERROR: Error processing PPTX {file_path}: {str(e)}")
            raise
    
    # Unused anymore but maybe keep
    def _pptx_to_markdown(self, prs: 'Presentation') -> str:
        """Convert PPTX presentation to Markdown."""
        markdown_lines = []
        
        for idx, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"## Slide {idx}\n")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown_lines.append(f"{shape.text.strip()}\n")
            
            markdown_lines.append("\n---\n")
        
        return "\n".join(markdown_lines)
    
    # Unused anymore but maybe keep
    def _pptx_to_pdf(self, prs: 'Presentation', stem: str):
        """Convert PPTX presentation to PDF using ReportLab (fallback only).
        Uses safe stem to avoid Windows path length issues."""
        if not REPORTLAB_AVAILABLE:
            print("WARNING: ReportLab not available, cannot convert PPTX to PDF")
            return
        
        # Use safe stem
        pdf_path = self.pdf_dir / f"{stem}.pdf"
        
        try:
            pdf_doc = SimpleDocTemplate(
                str(pdf_path),
                pagesize=A4 if self.config.pdf_page_size == "A4" else letter
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            for idx, slide in enumerate(prs.slides, 1):
                # Add slide header
                story.append(Paragraph(f"Slide {idx}", styles['Heading1']))
                story.append(Spacer(1, 0.3 * inch))
                
                # Add text content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text.strip(), styles['Normal']))
                        story.append(Spacer(1, 0.2 * inch))
                
                # Page break between slides
                if idx < len(prs.slides):
                    story.append(PageBreak())
            
            pdf_doc.build(story)
            print(f"Created PDF: {pdf_path}")
        
        except Exception as e:
            print(f"ERROR: Error creating PDF from PPTX: {str(e)}")
    
    # ======================== HTML Normalization ========================
    
    def _normalize_html(self, file_path: Path, stem: str):
        """Normalize HTML to Markdown and PDF."""
        if not BS4_AVAILABLE:
            print("WARNING: BeautifulSoup not available, skipping HTML processing")
            return
        
        try:
            # Convert to PDF using LibreOffice (preserves layout, images, styling)
            if self.config.generate_pdf:
                print("  → Trying LibreOffice for HTML to PDF conversion...")
                if self._html_to_pdf_libreoffice(file_path, stem):
                    print("  → ✓ LibreOffice HTML→PDF conversion successful")
                else:
                    print("  → LibreOffice not available for HTML, PDF not generated")
        
        except Exception as e:
            print(f"ERROR: Error processing HTML {file_path}: {str(e)}")
            raise
    
    def _html_to_pdf_libreoffice(self, file_path: Path, stem: str) -> bool:
        """Convert HTML to PDF (Lambda first when enabled, then local LibreOffice)."""
        # Use safe stem to avoid Windows path length issues
        pdf_path = self.pdf_dir / f"{stem}.pdf"

        if self._should_use_lambda_office_pdf_conversion():
            print("    → Trying AWS Lambda LibreOffice conversion for HTML...")
            if self._office_to_pdf_lambda(file_path, pdf_path):
                return True
            print("    ✗ AWS Lambda HTML conversion failed; trying local LibreOffice...")

        return self._office_to_pdf_local(file_path, pdf_path)
    
    # ======================== Excel Normalization ========================
    
    def _normalize_excel(self, file_path: Path, stem: str):
        """Parse Excel files using custom XML-based parser (xlsx_reader_v2).
        
        Produces structured JSON in excel_parsed/ that preserves:
        - Table structures with proper headers
        - Charts with data series
        - Shapes and embedded images
        - Merged cells, hyperlinks, styles
        
        The JSON is consumed later by _run_excel_processing() in the pipeline
        to produce RAG-ready chunks via ExcelPreprocessor.
        """
        ext = file_path.suffix.lower()
        
        # Only our custom parser handles .xlsx and .xlsm natively
        if ext in ('.xlsx', '.xlsm'):
            try:
                from .xlsx_reader_v2 import process_excel_file
                
                json_output = process_excel_file(
                    excel_path=file_path,
                    output_dir=self.excel_parsed_dir,
                    parsed_parent=self.excel_parsed_dir / "_parsed",
                )
                print(f"  → ✓ Excel parsed to JSON: {json_output.name}")
                
            except Exception as e:
                print(f"  → ✗ Custom Excel parser failed: {e}")
                print(f"  → Falling back to Docling processing")
                # Keep the copy in originals_dir so Docling can try
        
        elif ext == '.xls':
            # .xls (legacy binary format) — convert to .xlsx via LibreOffice then parse
            try:
                from .xls_reader import XlsParser
                from .xlsx_reader_v2 import process_excel_file

                xlsx_path = XlsParser().convert_xls_to_xlsx(str(file_path))
                json_output = process_excel_file(
                    excel_path=Path(xlsx_path),
                    output_dir=self.excel_parsed_dir,
                    parsed_parent=self.excel_parsed_dir / "_parsed",
                )
                print(f"  → ✓ .xls converted and parsed to JSON: {json_output.name}")
            except Exception as e:
                print(f"  → ✗ .xls custom parsing failed: {e}")
                print(f"  → Falling back to Docling processing")
    
    # ======================== CSV Normalization ========================
    
    def _normalize_csv(self, file_path: Path, stem: str):
        """Normalize CSV files - no conversion needed, Docling will handle it."""
        # CSV files are already copied to original_files for Docling processing
        # No markdown conversion needed here
        print(f"  → CSV file ready for Docling processing")
        pass
    
    def _copy_to_originals_only(self, file_path: Path, stem: str):
        """Copy files to original_files only (already done in _normalize_file)."""
        # File already copied to original_files for Docling processing
        print(f"  → File ready for Docling processing")
        pass
        
    # ======================== Image Normalization ========================
    
    def _normalize_image(self, file_path: Path, stem: str):
        """Normalize image files - NO PDF conversion, let Docling handle natively."""
        # CRITICAL FIX: Do NOT convert images to PDF!
        # Docling processes images better than converted PDFs
        # Images already copied to original_files/, no additional processing needed
        
        print(f"  → Image ready for Docling native processing (no PDF conversion)")
        
        # Note: Image was already copied to original_files/ in _normalize_file()
        # Docling will process it directly with VLM, picture classification, etc.
    
    # ======================== Direct Copy Operations ========================
    
    def _copy_pdf(self, file_path: Path, stem: str):
        """Classify and process PDF files.

        Born-digital PDFs are parsed with pdf_reader.py (fast, no OCR) and
        produce a heading-tree JSON consumed by the pipeline's pdf_preprocessor.
        Scanned/hybrid PDFs are copied to normalized_pdfs/ for Docling Stage 3.
        """
        # Always copy to normalized_pdfs/ for Stage-3 Docling compatibility
        if self.config.generate_pdf:
            dest = self.pdf_dir / f"{stem}.pdf"
            shutil.copy(file_path, dest)

        # --- Classify ---
        try:
            from .pdf_classifier import PdfClassifier, PdfType
            classifier = PdfClassifier()
            classification = classifier.classify(file_path)

            print(f"  PDF classified as: {classification.pdf_type.value} "
                  f"(confidence={classification.confidence:.2f}, "
                  f"pages={classification.page_count})")

            import json as _json
            meta_path = self.metadata_dir / f"{stem}_pdf_classification.json"
            meta = {
                "pdf_type": classification.pdf_type.value,
                "confidence": classification.confidence,
                "pdf_version": classification.pdf_version,
                "has_structure_tree": classification.has_structure_tree,
                "page_count": classification.page_count,
                "signals": classification.signals,
            }
            with open(meta_path, "w", encoding="utf-8") as f:
                _json.dump(meta, f, ensure_ascii=False, indent=2, default=str)

            # Born-digital → parse with pdf_reader (fast path, no OCR)
            if classification.pdf_type == PdfType.BORN_DIGITAL:
                self._parse_pdf_with_reader(file_path, stem)
            else:
                print(f"  → Scanned/hybrid PDF — leaving in normalized_pdfs/ for Docling Stage 3")

        except Exception as e:
            print(f"  WARNING: PDF classification failed ({e}), "
                  "file left in normalized_pdfs/ for Stage-3 Docling")

    def _parse_pdf_with_reader(self, file_path: Path, stem: str):
        """Parse a born-digital PDF into heading-tree JSON using pdf_reader.py."""
        import json as _json
        from .pdf_reader import CustomPdfConfig, CustomPdfReader

        parsed_output_dir = self.pdf_parsed_dir / "_parsed" / stem
        parsed_output_dir.mkdir(parents=True, exist_ok=True)

        try:
            content_source = (
                os.getenv("PDF_READER_CONTENT_SOURCE")
                or self.config.pdf_content_source
                or "hybrid"
            ).strip().lower()
            if content_source not in {"pymupdf", "docling", "hybrid"}:
                print(
                    f"  WARNING: invalid PDF_READER_CONTENT_SOURCE={content_source!r}; "
                    "falling back to hybrid"
                )
                content_source = "hybrid"

            enable_ocr = bool(self.config.pdf_reader_enable_ocr)
            extract_images = bool(self.config.pdf_reader_extract_images)
            print(
                "  → PDF content extraction: "
                f"{content_source} (ocr={enable_ocr}, images={extract_images})"
            )

            reader = CustomPdfReader(
                CustomPdfConfig(
                    enable_ocr=enable_ocr,
                    extract_images=extract_images,
                    extract_tables=False,
                    content_source=content_source,
                    runtime_yaml=self.config.runtime_yaml,
                )
            )
            tree = reader.read(
                str(file_path),
                output_dir=str(parsed_output_dir),
                skip_ocr=not enable_ocr,
            )

            out_json = self.pdf_parsed_dir / f"{stem}.json"
            with open(out_json, "w", encoding="utf-8") as f:
                _json.dump(tree, f, ensure_ascii=False, indent=2)

            print(f"  → ✓ Born-digital PDF parsed: {out_json.name} "
                  f"({len(tree)} top-level sections)")
        except Exception as e:
            print(f"  WARNING: pdf_reader failed ({e}), "
                  "file left in normalized_pdfs/ for Stage-3 Docling")
    
    def _txt_to_md(self, file_path: Path, stem: str):
        """Convert TXT to MD format with safe name."""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Use safe stem
        md_path = self.markdown_dir / f"{stem}.md"
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        print(f"  → Converted TXT to MD: {md_path.name}")
    
    def _copy_text(self, file_path: Path, stem: str):
        """Copy markdown files to normalized directory with safe name."""
        if self.config.generate_markdown:
            dest = self.markdown_dir / f"{stem}.md"
            shutil.copy(file_path, dest)
            print(f"  → Copied markdown file: {dest.name}")
    
    # ======================== Helper Methods ========================
    
    def _save_markdown(self, stem: str, content: str):
        """Save markdown content to file."""
        md_path = self.markdown_dir / f"{stem}.md"
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(content)
        print(f"Created markdown: {md_path}")
    
    def _save_statistics(self):
        """Save processing statistics to JSON."""
        stats_path = self.metadata_dir / "normalization_stats.json"
        with open(stats_path, 'w', encoding='utf-8') as f:
            json.dump(self.stats, f, indent=2)
        print(f"Saved statistics to: {stats_path}")
    
    def __enter__(self):
        """Context manager entry."""
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit."""
        if exc_type is not None:
            print(f"ERROR: Error during normalization: {exc_val}")
        return False


# ======================== Utility Functions ========================

def normalize_documents(
    input_dir: Union[str, Path],
    output_dir: Union[str, Path],
    config: Optional[NormalizerConfig] = None
) -> Dict:
    """
    Convenience function to normalize documents.
    
    Args:
        input_dir: Directory containing source documents
        output_dir: Directory for normalized outputs
        config: Optional normalization configuration
    
    Returns:
        Dictionary with processing statistics
    """
    with DocumentNormalizer(input_dir, output_dir, config) as normalizer:
        return normalizer.normalize_batch()


if __name__ == "__main__":
    # Example usage
    config = NormalizerConfig(
        generate_pdf=True,
        generate_markdown=True,
        excel_to_markdown=True,
        csv_to_markdown=True
    )
    
    stats = normalize_documents(
        input_dir="input",
        output_dir="output/normalized",
        config=config
    )
    
    print(f"\nNormalization Summary:")
    print(f"Total files: {stats['total_files']}")
    print(f"Successfully normalized: {stats['normalized_files']}")
    print(f"Failed: {stats['failed_files']}")
    print(f"\nBy type: {stats['by_type']}")
