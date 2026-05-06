#!/usr/bin/env python3
"""Generate pseudo section-hierarchy GT for Stanford/CS PDF and PPTX files.

The output is intentionally separate from OmniDocBench-style metrics. It uses
course folder structure as the upper hierarchy and native document signals
(PDF outline/page titles, PPTX slide titles) as document-local sections.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from statistics import median
from typing import Any, Dict, Iterable, List, Optional, Sequence

BACKEND_ROOT = Path(__file__).resolve().parents[3]
REPO_ROOT = BACKEND_ROOT.parent

DEFAULT_ROOTS = [
    "backend/input/CS131_notes_zh-CN-master",
    "backend/input/CS193P-Fall-2017-DEMO-master",
    "backend/input/CS228_PGM-master",
    "backend/input/CS231n-2017-master",
    "backend/input/Stanford-University-Algorithms-Design-and-Analysis-master",
    "backend/input/Stanford-University-Statistical-Learning-master",
    "backend/input/cs229-2018-autumn-main",
    "backend/input/stanford-cs161-master",
]


@dataclass
class SectionNode:
    id: str
    title: str
    level: int
    parent_id: Optional[str]
    source: str
    page_start: Optional[int] = None
    slide_start: Optional[int] = None
    confidence: float = 1.0
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentHierarchy:
    doc_id: str
    source_path: str
    file_type: str
    course: str
    module_path: List[str]
    extraction_method: str
    page_count: Optional[int] = None
    slide_count: Optional[int] = None
    nodes: List[SectionNode] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--roots",
        nargs="*",
        default=DEFAULT_ROOTS,
        help="Course roots to scan. Defaults to Stanford/CS folders under backend/input.",
    )
    parser.add_argument(
        "--output-dir",
        default=str(BACKEND_ROOT / "output" / "evaluation" / "stanford_section_hierarchy_pdf_pptx"),
    )
    parser.add_argument("--max-files", type=int, default=None)
    parser.add_argument("--max-pages-per-pdf", type=int, default=None)
    parser.add_argument("--include-pdf-outlines", action="store_true", default=True)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    roots = [Path(root).resolve() for root in args.roots]
    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    files = list(discover_files(roots))
    if args.max_files:
        files = files[: args.max_files]

    documents: List[DocumentHierarchy] = []
    failures: List[Dict[str, str]] = []
    for idx, path in enumerate(files, 1):
        try:
            hierarchy = build_document_hierarchy(path, roots, max_pages_per_pdf=args.max_pages_per_pdf)
            documents.append(hierarchy)
        except Exception as exc:
            failures.append({"path": str(path), "error": f"{type(exc).__name__}: {exc}"})
        if idx % 25 == 0:
            print(f"processed {idx}/{len(files)} files; documents={len(documents)} failures={len(failures)}", flush=True)

    payload = {
        "metadata": {
            "schema": "stanford_pdf_pptx_section_hierarchy_pseudo_gt",
            "version": "1.0",
            "roots": [str(root) for root in roots],
            "file_count": len(files),
            "document_count": len(documents),
            "failure_count": len(failures),
            "modalities": ["pdf", "pptx"],
            "method": "folder_structure_plus_pdf_outline_or_page_titles_plus_pptx_slide_titles",
        },
        "documents": [document_to_dict(document) for document in documents],
        "failures": failures,
    }
    summary = build_summary(documents, failures, files)

    hierarchy_path = out_dir / "section_hierarchy_pseudo_gt.json"
    summary_path = out_dir / "summary.json"
    failures_path = out_dir / "failures.json"
    hierarchy_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    summary_path.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    failures_path.write_text(json.dumps(failures, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"hierarchy": str(hierarchy_path), "summary": str(summary_path), "failures": str(failures_path)}, indent=2))
    return 0


def discover_files(roots: Sequence[Path]) -> Iterable[Path]:
    for root in roots:
        if not root.exists():
            continue
        candidates = [path for path in root.rglob("*") if path.is_file() and path.suffix.lower() in {".pdf", ".pptx"}]
        for path in sorted(candidates, key=lambda p: str(p).lower()):
            if path.name.startswith("~$"):
                continue
            yield path


def build_document_hierarchy(path: Path, roots: Sequence[Path], *, max_pages_per_pdf: Optional[int]) -> DocumentHierarchy:
    course_root = next((root for root in roots if path.is_relative_to(root)), path.parent)
    rel = path.relative_to(course_root)
    course = course_root.name
    module_path = [clean_title(part) for part in rel.parts[:-1]]
    doc_title = clean_title(path.stem)
    doc_id = stable_id(str(path.resolve()))
    file_type = path.suffix.lower().lstrip(".")
    hierarchy = DocumentHierarchy(
        doc_id=doc_id,
        source_path=str(path),
        file_type=file_type,
        course=course,
        module_path=module_path,
        extraction_method="",
    )

    parent_id: Optional[str] = None
    level = 1
    course_id = stable_id(f"{doc_id}:course:{course}")
    hierarchy.nodes.append(SectionNode(course_id, course, level, None, "folder", confidence=1.0))
    parent_id = course_id
    for module in module_path:
        level += 1
        module_id = stable_id(f"{doc_id}:module:{level}:{module}")
        hierarchy.nodes.append(SectionNode(module_id, module, level, parent_id, "folder", confidence=1.0))
        parent_id = module_id

    level += 1
    document_node_id = stable_id(f"{doc_id}:document:{doc_title}")
    hierarchy.nodes.append(SectionNode(document_node_id, doc_title, level, parent_id, "filename", confidence=0.95))

    if file_type == "pdf":
        add_pdf_sections(hierarchy, path, parent_id=document_node_id, level=level + 1, max_pages=max_pages_per_pdf)
    elif file_type == "pptx":
        add_pptx_sections(hierarchy, path, parent_id=document_node_id, level=level + 1)
    else:
        hierarchy.warnings.append(f"unsupported file type: {file_type}")
    return hierarchy


def add_pdf_sections(hierarchy: DocumentHierarchy, path: Path, *, parent_id: str, level: int, max_pages: Optional[int]) -> None:
    import fitz

    doc = fitz.open(str(path))
    hierarchy.page_count = doc.page_count
    toc = doc.get_toc(simple=True)
    if toc:
        hierarchy.extraction_method = "pdf_outline"
        stack: Dict[int, str] = {level - 1: parent_id}
        for idx, (toc_level, title, page_no) in enumerate(toc, 1):
            cleaned = clean_title(title)
            if not cleaned:
                continue
            node_level = level + max(0, int(toc_level) - 1)
            node_id = stable_id(f"{hierarchy.doc_id}:pdf_outline:{idx}:{cleaned}:{page_no}")
            node_parent = stack.get(node_level - 1, parent_id)
            hierarchy.nodes.append(
                SectionNode(
                    node_id,
                    cleaned,
                    node_level,
                    node_parent,
                    "pdf_outline",
                    page_start=int(page_no) if page_no else None,
                    confidence=0.98,
                )
            )
            stack[node_level] = node_id
        return

    hierarchy.extraction_method = "pdf_page_title"
    page_limit = min(doc.page_count, max_pages) if max_pages else doc.page_count
    previous_title = ""
    for page_index in range(page_limit):
        page = doc.load_page(page_index)
        title = extract_pdf_page_title(page)
        if not title or title == previous_title:
            continue
        previous_title = title
        node_id = stable_id(f"{hierarchy.doc_id}:pdf_page:{page_index + 1}:{title}")
        hierarchy.nodes.append(
            SectionNode(
                node_id,
                title,
                level,
                parent_id,
                "pdf_page_title",
                page_start=page_index + 1,
                confidence=0.72,
            )
        )
    if max_pages and doc.page_count > max_pages:
        hierarchy.warnings.append(f"pdf truncated at {max_pages}/{doc.page_count} pages")


def extract_pdf_page_title(page: Any) -> str:
    data = page.get_text("dict")
    page_height = float(page.rect.height or 1.0)
    lines: List[Dict[str, Any]] = []
    for block in data.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            text = clean_title(" ".join(str(span.get("text") or "") for span in spans))
            if not useful_title(text):
                continue
            sizes = [float(span.get("size") or 0.0) for span in spans]
            bbox = line.get("bbox") or [0, page_height, 0, page_height]
            y0 = float(bbox[1])
            lines.append({"text": text, "size": max(sizes) if sizes else 0.0, "y0": y0})
    if not lines:
        return ""
    top_lines = [line for line in lines if line["y0"] <= page_height * 0.42]
    pool = top_lines or lines[:8]
    sizes = [line["size"] for line in pool if line["size"] > 0]
    threshold = median(sizes) if sizes else 0.0
    pool.sort(key=lambda item: (-item["size"], item["y0"], len(item["text"])))
    for item in pool:
        if item["size"] >= threshold and useful_title(item["text"]):
            return item["text"]
    return ""


def add_pptx_sections(hierarchy: DocumentHierarchy, path: Path, *, parent_id: str, level: int) -> None:
    from pptx import Presentation

    presentation = Presentation(str(path))
    hierarchy.slide_count = len(presentation.slides)
    hierarchy.extraction_method = "pptx_slide_title"
    previous_title = ""
    for slide_index, slide in enumerate(presentation.slides, 1):
        title = extract_pptx_slide_title(slide)
        if not title or title == previous_title:
            continue
        previous_title = title
        node_id = stable_id(f"{hierarchy.doc_id}:pptx_slide:{slide_index}:{title}")
        hierarchy.nodes.append(
            SectionNode(
                node_id,
                title,
                level,
                parent_id,
                "pptx_slide_title",
                slide_start=slide_index,
                confidence=0.82,
            )
        )


def extract_pptx_slide_title(slide: Any) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        text = clean_title(getattr(title_shape, "text", "") or "")
        if useful_title(text):
            return text
    candidates: List[Dict[str, Any]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = clean_title(getattr(shape, "text", "") or "")
        if not useful_title(text):
            continue
        font_sizes = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    font_sizes.append(float(run.font.size.pt))
        candidates.append(
            {
                "text": text,
                "top": int(getattr(shape, "top", 0) or 0),
                "size": max(font_sizes) if font_sizes else 0.0,
            }
        )
    if not candidates:
        return ""
    candidates.sort(key=lambda item: (-item["size"], item["top"], len(item["text"])))
    return candidates[0]["text"]


def useful_title(text: str) -> bool:
    if not text:
        return False
    if len(text) < 3 or len(text) > 180:
        return False
    if re.fullmatch(r"[\d\s./:-]+", text):
        return False
    if text.lower() in {"contents", "overview", "outline"}:
        return True
    return bool(re.search(r"[A-Za-z\u4e00-\u9fff]", text))


def clean_title(value: str) -> str:
    value = value or ""
    value = value.replace("\u00a0", " ")
    value = re.sub(r"\s+", " ", value)
    value = value.strip(" \t\r\n#*-:|")
    return value.strip()


def stable_id(value: str) -> str:
    return hashlib.sha1(value.encode("utf-8")).hexdigest()[:16]


def document_to_dict(document: DocumentHierarchy) -> Dict[str, Any]:
    data = asdict(document)
    return data


def build_summary(documents: Sequence[DocumentHierarchy], failures: Sequence[Dict[str, str]], files: Sequence[Path]) -> Dict[str, Any]:
    by_type: Dict[str, int] = {}
    by_method: Dict[str, int] = {}
    node_counts: List[int] = []
    for document in documents:
        by_type[document.file_type] = by_type.get(document.file_type, 0) + 1
        by_method[document.extraction_method] = by_method.get(document.extraction_method, 0) + 1
        node_counts.append(len(document.nodes))
    return {
        "file_count": len(files),
        "document_count": len(documents),
        "failure_count": len(failures),
        "by_type": by_type,
        "by_extraction_method": by_method,
        "total_nodes": sum(node_counts),
        "avg_nodes_per_document": (sum(node_counts) / len(node_counts)) if node_counts else 0.0,
        "failures_preview": list(failures[:20]),
    }


if __name__ == "__main__":
    raise SystemExit(main())
